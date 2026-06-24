#!/usr/bin/env python3
"""Simple PDF-to-editable-PPTX pipeline.

This is the default small-tool flow for notebooklm-pdf-to-ppt:
PDF -> rendered page PNGs -> OCR lines -> locally cleaned background ->
python-pptx editable text overlay.

It intentionally avoids the older multi-model fusion stack. The goal is to
make OCR quality and rebuild quality easy to inspect independently.
"""

from __future__ import annotations

import argparse
import json
import os
import re
import subprocess
import sys
from pathlib import Path
from typing import Any



EMU_PER_INCH = 914400
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
DEFAULT_CJK_FONT = "Noto Sans SC"
DEFAULT_LATIN_FONT = "Inter"
FALLBACK_LATIN_FONT = "Arial"
PLAYFUL_LATIN_FONTS = ("Comic Sans MS", "Chalkboard SE", "Marker Felt", DEFAULT_LATIN_FONT, FALLBACK_LATIN_FONT)
PLAYFUL_CJK_FONTS = ("ZCOOL KuaiLe", "Noto Sans SC", "Source Han Sans CN")
FONT_FILES = {
    "Noto Sans SC": (
        "~/Library/Fonts/NotoSansSC-Variable.ttf",
        "/Library/Fonts/NotoSansSC-Variable.ttf",
    ),
    "Source Han Sans CN": (
        "~/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "~/Library/Fonts/SourceHanSansCN-Bold.otf",
        "/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "/Library/Fonts/SourceHanSansCN-Bold.otf",
    ),
    "Arial": (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    ),
    "Times New Roman": (
        "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
        "/System/Library/Fonts/Supplemental/Times New Roman Bold.ttf",
    ),
    "Comic Sans MS": (
        "/System/Library/Fonts/Supplemental/Comic Sans MS.ttf",
        "/System/Library/Fonts/Supplemental/Comic Sans MS Bold.ttf",
    ),
    "Chalkboard SE": (
        "/System/Library/Fonts/Supplemental/ChalkboardSE.ttc",
    ),
}
CJK_FONT_FIT_CANDIDATES = ("Noto Sans SC", "Source Han Sans CN", "Arial")
LATIN_FONT_FIT_CANDIDATES = ("Arial", "Comic Sans MS", "Chalkboard SE", "Times New Roman")


def command_output(cmd: str, timeout: int = 5) -> str:
    try:
        result = subprocess.run(
            ["/bin/zsh", "-lc", cmd],
            check=True,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
    except (OSError, subprocess.SubprocessError):
        return ""
    return result.stdout.strip()


def font_available(font_name: str) -> bool:
    escaped = font_name.replace('"', '\\"')
    family = command_output(f'fc-match -f "%{{family}}" "{escaped}" 2>/dev/null')
    if family and font_name.lower() in {part.strip().lower() for part in family.split(",")}:
        return True
    # macOS commonly has these families even when fontconfig is unavailable.
    mac_fonts = {
        "arial",
        "comic sans ms",
        "chalkboard se",
        "marker felt",
        "times new roman",
        "noto sans sc",
        "source han sans cn",
    }
    return font_name.lower() in mac_fonts


def first_available_font(candidates: tuple[str, ...]) -> str:
    for font_name in candidates:
        if font_available(font_name):
            return font_name
    return candidates[-1]


def parse_pages(pages: str | None, total: int | None = None) -> list[int]:
    if not pages:
        if total is None:
            raise ValueError("--pages is required when total page count is unknown")
        return list(range(total))
    selected: set[int] = set()
    for part in pages.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start, end = int(start_s), int(end_s)
            selected.update(range(start - 1, end))
        else:
            selected.add(int(part) - 1)
    ordered = sorted(selected)
    if total is None:
        return [idx for idx in ordered if idx >= 0]
    return [idx for idx in ordered if 0 <= idx < total]


def contains_cjk(text: str) -> bool:
    return bool(re.search(r"[\u3400-\u9fff]", text))


def font_for_text(text: str) -> str:
    if contains_cjk(text):
        return DEFAULT_CJK_FONT
    return first_available_font((DEFAULT_LATIN_FONT, FALLBACK_LATIN_FONT))


def classify_text_role(item: dict[str, Any], page_w: int, page_h: int) -> str:
    y = float(item.get("y") or 0)
    h = float(item.get("height") or 0)
    w = float(item.get("width") or 0)
    if y < page_h * 0.24 and h > page_h * 0.07:
        return "title"
    if w > page_w * 0.30 and h > page_h * 0.035:
        return "body"
    return "label"


def apply_font_policy(item: dict[str, Any], page_w: int, page_h: int) -> None:
    role = classify_text_role(item, page_w, page_h)
    item["role"] = role
    text = str(item.get("text") or "")
    if role == "title":
        if contains_cjk(text):
            item["font_family"] = first_available_font(PLAYFUL_CJK_FONTS)
        else:
            item["font_family"] = first_available_font(PLAYFUL_LATIN_FONTS)
    else:
        item["font_family"] = font_for_text(text)


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"\s+([，。！？：；、])", r"\1", text)
    return text


def useful_text(text: str, conf: float, width: float, height: float) -> bool:
    if not text:
        return False
    visible = re.findall(r"[A-Za-z0-9\u3400-\u9fff]", text)
    if len(visible) < 2:
        return False
    if width < 12 or height < 8:
        return False
    if conf < 25 and len(visible) < 5:
        return False
    return True


def is_footer_brand(item: dict[str, Any], page_w: int, page_h: int) -> bool:
    text = str(item.get("text") or "").strip().lower()
    y = float(item.get("y") or 0)
    h = float(item.get("height") or 0)
    if y < page_h * 0.94:
        return False
    if h > page_h * 0.04:
        return False
    return text in {"notebooklm"} or len(text) <= 18


def overlap_ratio(a1: float, a2: float, b1: float, b2: float) -> float:
    overlap = max(0.0, min(a2, b2) - max(a1, b1))
    return overlap / max(1.0, min(a2 - a1, b2 - b1))


def can_merge_line(left: dict[str, Any], right: dict[str, Any]) -> bool:
    if left.get("font_family") != right.get("font_family"):
        return False
    ly1, ly2 = float(left["y"]), float(left["y"] + left["height"])
    ry1, ry2 = float(right["y"]), float(right["y"] + right["height"])
    if overlap_ratio(ly1, ly2, ry1, ry2) < 0.58:
        return False
    gap = float(right["x"]) - float(left["x"] + left["width"])
    max_h = max(float(left["height"]), float(right["height"]))
    if gap < 0:
        # OCR can over-expand neighboring word/phrase boxes on the same
        # baseline, especially for CJK text on textured or illustrated slides.
        # If the right box still advances in reading direction, treat the
        # negative gap as bbox overlap rather than a separate text object.
        advance = float(right["x"]) - float(left["x"])
        min_w = min(float(left["width"]), float(right["width"]))
        return gap >= -max(180.0, max_h * 2.5) and advance >= max(12.0, min_w * 0.24)
    return gap <= max(32.0, max_h * 1.25)


def merge_text_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    rows = sorted(items, key=lambda item: (float(item["y"]), float(item["x"])))
    merged: list[dict[str, Any]] = []
    for item in rows:
        current = dict(item)
        if merged and can_merge_line(merged[-1], current):
            prev = merged[-1]
            x1 = min(float(prev["x"]), float(current["x"]))
            y1 = min(float(prev["y"]), float(current["y"]))
            x2 = max(float(prev["x"] + prev["width"]), float(current["x"] + current["width"]))
            y2 = max(float(prev["y"] + prev["height"]), float(current["y"] + current["height"]))
            prev["text"] = clean_text(f"{prev['text']} {current['text']}")
            prev["x"] = round(x1, 2)
            prev["y"] = round(y1, 2)
            prev["width"] = round(x2 - x1, 2)
            prev["height"] = round(y2 - y1, 2)
            prev["confidence"] = round((float(prev.get("confidence", 0)) + float(current.get("confidence", 0))) / 2, 2)
            prev["font_size_px"] = round(max(float(prev.get("font_size_px", 0)), float(current.get("font_size_px", 0))), 2)
            prev.setdefault("word_boxes", []).extend(current.get("word_boxes") or [current])
        else:
            merged.append(current)
    return merged


def calibrate_font_size(item: dict[str, Any]) -> None:
    source = str(item.get("source") or "")
    height = float(item.get("height") or 0)
    current = float(item.get("font_size_px") or height)
    role = str(item.get("role") or "body")
    if source == "paddle":
        # The paddle worker now sizes from MEASURED ink height (style_probe),
        # so trust it and only clamp pathological extremes. The old 0.50/0.54
        # bbox factors halved every glyph and were the main "字号偏小" cause.
        if role == "title":
            max_factor = 1.18
        elif contains_cjk(str(item.get("text") or "")):
            max_factor = 0.95
        else:
            max_factor = 0.90
        item["font_size_px"] = round(max(6.0, min(current, height * max_factor)), 2)

def normalize_ocr_slide(page: dict[str, Any]) -> dict[str, Any]:
    page_w = int(page.get("width") or 0)
    page_h = int(page.get("height") or 0)
    texts = [dict(item) for item in page.get("texts", []) if not is_footer_brand(item, page_w, page_h)]
    for item in texts:
        apply_font_policy(item, page_w, page_h)
        calibrate_font_size(item)
    merged_texts = merge_text_items(texts)
    for item in merged_texts:
        apply_font_policy(item, page_w, page_h)
        calibrate_font_size(item)
    raw_mask_words = []
    for item in page.get("mask_words", page["texts"]):
        if is_footer_brand(item, page_w, page_h):
            continue
        mask_item = dict(item)
        apply_font_policy(mask_item, page_w, page_h)
        raw_mask_words.append(mask_item)
    page["texts"] = merged_texts
    # Background cleanup uses both raw OCR component boxes and merged editable
    # boxes. Merged boxes cover OCR misses between fragments, especially for
    # large titles. On illustrated backgrounds this can leave visible blocks;
    # use model-clean for high-fidelity repair.
    page["mask_words"] = raw_mask_words + merged_texts
    return page


_SECTION_NUMERAL_RE = re.compile(r"^[一二三四五六七八九十百0-9]{1,3}[、.，]?$")


def merge_section_numeral_prefix(page: dict[str, Any]) -> None:
    """Re-attach a lone section numeral ("四") to its heading -> "四、标题".

    PaddleOCR detects the numeral as its own box and often drops the "、". The
    numeral carries section meaning, so prepend it (with a restored "、") to the
    nearest same-row heading to its right instead of leaving it floating.
    """
    texts = page.get("texts", [])
    numerals = [
        item for item in texts
        if _SECTION_NUMERAL_RE.match(str(item.get("text") or "").strip())
        and not contains_cjk(re.sub(r"[一二三四五六七八九十百]", "", str(item.get("text") or "")))
    ]
    for num in numerals:
        ny = float(num.get("y") or 0)
        nh = float(num.get("height") or 0)
        nx2 = float(num.get("x") or 0) + float(num.get("width") or 0)
        target = None
        best_dx = 10**9
        for item in texts:
            if item is num:
                continue
            iy = float(item.get("y") or 0)
            ih = float(item.get("height") or 0)
            # same row: vertical centers overlap
            if abs((iy + ih / 2) - (ny + nh / 2)) > max(nh, ih) * 0.6:
                continue
            ix = float(item.get("x") or 0)
            dx = ix - nx2
            if dx < -nh or dx > max(nh, ih) * 4:  # must sit just to the right
                continue
            if dx < best_dx:
                best_dx = dx
                target = item
        if target is None:
            continue
        digits = re.sub(r"[、.，]", "", str(num.get("text") or "").strip())
        target_text = str(target.get("text") or "").strip()
        old_x = float(target.get("x") or 0)
        old_right = old_x + float(target.get("width") or 0)
        new_x = min(old_x, float(num.get("x") or 0))
        target["text"] = f"{digits}、{target_text}"
        target["x"] = round(new_x, 2)
        target["width"] = round(old_right - new_x, 2)
        target["sectionNumeralRestored"] = digits
    if numerals:
        page["texts"] = [item for item in texts if item not in numerals]


def can_merge_cjk_short_continuation(prev: dict[str, Any], current: dict[str, Any], page_w: int, page_h: int) -> bool:
    prev_text = str(prev.get("text") or "").strip()
    current_text = str(current.get("text") or "").strip()
    if not prev_text or not current_text:
        return False
    if str(prev.get("role") or "") == "title" or str(current.get("role") or "") == "title":
        return False
    if prev.get("styleGroup") or current.get("styleGroup") or prev.get("tableRow") or current.get("tableRow"):
        return False
    if not contains_cjk(prev_text) or not re.search(r"[\u3400-\u9fff]", current_text):
        return False
    current_visible = re.findall(r"[A-Za-z0-9\u3400-\u9fff]", current_text)
    if len(current_visible) > 3:
        return False
    if re.search(r"[。！？.!?…]$", prev_text):
        return False
    if re.match(r"^[QA][：:]\s+", current_text) or list_like_text(current_text):
        return False
    px = float(prev.get("x") or 0)
    py = float(prev.get("y") or 0)
    pw = float(prev.get("width") or 0)
    ph = float(prev.get("height") or 0)
    cx = float(current.get("x") or 0)
    cy = float(current.get("y") or 0)
    cw = float(current.get("width") or 0)
    ch = float(current.get("height") or 0)
    if cy < py:
        return False
    vertical_gap = cy - (py + ph)
    if vertical_gap > max(16.0, max(ph, ch) * 0.35):
        return False
    same_left = abs(cx - px) <= max(36.0, page_w * 0.022)
    centered_tail = abs((cx + cw / 2) - (px + pw / 2)) <= max(90.0, page_w * 0.05)
    overlap = max(0.0, min(px + pw, cx + cw) - max(px, cx)) / max(1.0, min(pw, cw))
    return same_left or centered_tail or overlap >= 0.55


def merge_cjk_short_continuations(page: dict[str, Any]) -> None:
    page_w = int(page.get("width") or 0)
    page_h = int(page.get("height") or 0)
    items = [dict(item) for item in sorted(page.get("texts", []), key=lambda item: (float(item.get("y") or 0), float(item.get("x") or 0)))]
    consumed: set[int] = set()
    repairs: list[dict[str, Any]] = []
    for idx, current in enumerate(items):
        if idx in consumed:
            continue
        best_idx: int | None = None
        best_score = 10**9
        cy = float(current.get("y") or 0)
        for prev_idx, prev in enumerate(items[:idx]):
            if prev_idx in consumed:
                continue
            if not can_merge_cjk_short_continuation(prev, current, page_w, page_h):
                continue
            score = abs(float(current.get("x") or 0) - float(prev.get("x") or 0)) + abs(cy - float(prev.get("y") or 0)) * 0.2
            if score < best_score:
                best_score = score
                best_idx = prev_idx
        if best_idx is None:
            continue
        prev = items[best_idx]
        old_text = str(prev.get("text") or "")
        x1 = min(float(prev.get("x") or 0), float(current.get("x") or 0))
        y1 = min(float(prev.get("y") or 0), float(current.get("y") or 0))
        x2 = max(float(prev.get("x") or 0) + float(prev.get("width") or 0), float(current.get("x") or 0) + float(current.get("width") or 0))
        y2 = max(float(prev.get("y") or 0) + float(prev.get("height") or 0), float(current.get("y") or 0) + float(current.get("height") or 0))
        prev["text"] = f"{clean_text(old_text)}\n{clean_text(str(current.get('text') or ''))}"
        prev["x"] = round(x1, 2)
        prev["y"] = round(y1, 2)
        prev["width"] = round(x2 - x1, 2)
        prev["height"] = round(y2 - y1, 2)
        prev["source"] = "cjk_short_continuation"
        prev["textSource"] = "ocr_cjk_short_continuation"
        prev["lineBreakSource"] = "ocr_visible_rows"
        prev["word_wrap"] = True
        prev.setdefault("word_boxes", []).extend(current.get("word_boxes") or [current])
        consumed.add(idx)
        repairs.append(
            {
                "type": "cjk_short_continuation_merge",
                "mergedText": prev["text"],
                "texts": [old_text, current.get("text")],
                "reason": "OCR split a short CJK tail onto the next visual row inside the same text container; merge into one editable text box while preserving the original visible row break.",
            }
        )
    if consumed:
        page["texts"] = [item for idx, item in enumerate(items) if idx not in consumed]
    if repairs:
        page.setdefault("layout_repairs", []).extend(repairs)


def merge_primary_top_band_heading(page: dict[str, Any]) -> None:
    page_h = int(page.get("height") or 0)
    if not page_h:
        return
    top_limit = page_h * 0.15
    top_items = [
        item for item in page.get("texts", [])
        if float(item.get("y") or 0) <= top_limit and float(item.get("height") or 0) >= page_h * 0.045
    ]
    if len(top_items) < 2:
        return
    top_text = " ".join(str(item.get("text") or "") for item in top_items)
    if "Key" not in top_text and not re.search(r"[一二三四五六七八九十]+、", top_text):
        return
    # Merge a heading's OWN fragments only. Anchor on the fragment that carries
    # the heading signal (a section numeral "X、", else "Key"), then extend to
    # neighbours that are both horizontally contiguous (small gap) AND similar
    # in height. This excludes decorative signs / vocab cards that happen to sit
    # in the top band (e.g. a "TINY SIPS" shop sign left of the title). A numeral
    # marks the heading START, so a numeral anchor never absorbs anything to its
    # left.
    page_w = int(page.get("width") or 0)
    gap_limit = max(140.0, page_w * 0.12)
    ordered_all = sorted(top_items, key=lambda item: float(item.get("x") or 0))
    numeral_re = re.compile(r"[一二三四五六七八九十]+、")
    anchor_idx = next((i for i, it in enumerate(ordered_all) if numeral_re.search(str(it.get("text") or ""))), None)
    anchor_is_numeral = anchor_idx is not None
    if anchor_idx is None:
        anchor_idx = next((i for i, it in enumerate(ordered_all) if "Key" in str(it.get("text") or "")), None)
    if anchor_idx is None:
        return
    anchor_h = float(ordered_all[anchor_idx].get("height") or 0)

    def _compatible(a: dict[str, Any], b: dict[str, Any]) -> bool:
        a_right = float(a.get("x") or 0) + float(a.get("width") or 0)
        if float(b.get("x") or 0) - a_right > gap_limit:
            return False
        bh = float(b.get("height") or 0)
        if anchor_h and bh and abs(bh - anchor_h) > max(40.0, anchor_h * 0.5):
            return False
        return True

    ordered = [ordered_all[anchor_idx]]
    j = anchor_idx
    while j + 1 < len(ordered_all) and _compatible(ordered_all[j], ordered_all[j + 1]):
        ordered.append(ordered_all[j + 1])
        j += 1
    if not anchor_is_numeral:
        k = anchor_idx
        while k - 1 >= 0 and _compatible(ordered_all[k - 1], ordered_all[k]):
            ordered.insert(0, ordered_all[k - 1])
            k -= 1
    if len(ordered) < 2:
        return
    x1 = min(float(item.get("x") or 0) for item in ordered)
    y1 = min(float(item.get("y") or 0) for item in ordered)
    x2 = max(float(item.get("x") or 0) + float(item.get("width") or 0) for item in ordered)
    y2 = max(float(item.get("y") or 0) + float(item.get("height") or 0) for item in ordered)
    text = clean_text(" ".join(str(item.get("text") or "") for item in ordered))
    merged = dict(ordered[0])
    merged.update(
        {
            "text": text,
            "x": round(x1, 2),
            "y": round(y1, 2),
            "width": round(x2 - x1, 2),
            "height": round(y2 - y1, 2),
            "confidence": round(sum(float(item.get("confidence") or 0) for item in ordered) / len(ordered), 2),
            "source": "primary_top_band_merge",
            "textSource": "ocr_primary_group",
            "word_boxes": ordered,
        }
    )
    apply_font_policy(merged, int(page.get("width") or 0), page_h)
    calibrate_font_size(merged)
    page["texts"] = sorted([item for item in page.get("texts", []) if item not in ordered] + [merged], key=lambda item: (float(item.get("y") or 0), float(item.get("x") or 0)))
    page.setdefault("layout_repairs", []).append(
        {
            "type": "top_band_primary_heading_merge",
            "mergedText": text,
            "mergedCount": len(ordered),
            "reason": "Primary OCR split a top heading into adjacent title fragments; merge by x-order before typography and rebuild.",
        }
    )


def paragraph_merge_candidate(item: dict[str, Any], page_w: int, page_h: int) -> bool:
    text = str(item.get("text") or "").strip()
    if not text:
        return False
    visible = re.findall(r"[A-Za-z0-9\u3400-\u9fff]", text)
    if len(visible) < 4:
        return False
    if str(item.get("role") or "") == "title":
        return False
    if item.get("styleGroup"):
        return False
    if item.get("tableRow"):
        # A glossary/table row (has a same-row partner in another column) must
        # not be vertically merged with the row above/below it.
        return False
    if list_like_text(text):
        # Q:/A: prompts can be multi-line paragraphs; glossary/list items should stay separate.
        if not re.match(r"^[QA][：:]\s+", text):
            return False
    if float(item.get("y") or 0) < page_h * 0.16:
        return False
    if float(item.get("width") or 0) < page_w * 0.045:
        return False
    return True


def same_paragraph_style(a: dict[str, Any], b: dict[str, Any]) -> bool:
    if a.get("font_family") != b.get("font_family"):
        return False
    # Compare bbox HEIGHT, not the derived font_size_px. Font size is now
    # estimated from measured ink height, so sibling lines of one sentence can
    # diverge a lot (ascenders/descenders) even though their boxes are nearly
    # the same height -- which was wrongly splitting sentences. The bbox height
    # is the stable signal; the merged paragraph is locked to a group size after.
    a_h = float(a.get("height") or 0)
    b_h = float(b.get("height") or 0)
    if not a_h or not b_h:
        return True
    return abs(a_h - b_h) <= max(14.0, min(a_h, b_h) * 0.45)


def can_merge_paragraph_line(prev: dict[str, Any], current: dict[str, Any], page_w: int, page_h: int) -> bool:
    if not paragraph_merge_candidate(prev, page_w, page_h) or not paragraph_merge_candidate(current, page_w, page_h):
        return False
    if not same_paragraph_style(prev, current):
        return False
    px = float(prev.get("x") or 0)
    py = float(prev.get("y") or 0)
    pw = float(prev.get("width") or 0)
    ph = float(prev.get("height") or 0)
    cx = float(current.get("x") or 0)
    cy = float(current.get("y") or 0)
    cw = float(current.get("width") or 0)
    ch = float(current.get("height") or 0)
    vertical_gap = cy - (py + ph)
    if vertical_gap < -max(ph, ch) * 0.25:
        return False
    if vertical_gap > max(18.0, max(ph, ch) * 0.85):
        return False
    same_left = abs(cx - px) <= max(42.0, page_w * 0.025)
    same_center = abs((cx + cw / 2) - (px + pw / 2)) <= max(55.0, page_w * 0.035)
    enough_horizontal_overlap = max(0.0, min(px + pw, cx + cw) - max(px, cx)) / max(1.0, min(pw, cw)) >= 0.45
    if not (same_left or same_center or enough_horizontal_overlap):
        return False
    prev_text = str(prev.get("text") or "").strip()
    current_text = str(current.get("text") or "").strip()
    if re.match(r"^[QA][：:]\s+", current_text):
        return False
    if re.match(r"^[QA][：:]\s+", prev_text):
        return True
    prev_terminal = bool(re.search(r"[.!?。！？…]$", prev_text))
    if prev_terminal and contains_cjk(prev_text) == contains_cjk(current_text):
        return False
    return True


def paragraph_line_ref(item: dict[str, Any]) -> dict[str, Any]:
    return {
        "text": item.get("text"),
        "x": item.get("x"),
        "y": item.get("y"),
        "width": item.get("width"),
        "height": item.get("height"),
        "font_size_px": item.get("font_size_px"),
        "font_family": item.get("font_family"),
        "source": item.get("source"),
    }


def _flow_join(parts: list[str]) -> str:
    """Join sentence fragments into one flowing string instead of hard breaks.

    OCR splits a sentence into visual rows; forcing those rows back with "\n"
    produces a rigid break that no longer matches the resized text. Joining them
    (a space at Latin word boundaries, nothing between CJK) lets word_wrap reflow
    the sentence as one segment.
    """
    out = ""
    for seg in parts:
        seg = str(seg or "").strip()
        if not seg:
            continue
        if not out:
            out = seg
        elif contains_cjk(out[-1]) or contains_cjk(seg[0]):
            out += seg
        else:
            out += " " + seg
    return out


def merge_paragraph_group(lines: list[dict[str, Any]]) -> dict[str, Any]:
    ordered = sorted(lines, key=lambda item: (float(item.get("y") or 0), float(item.get("x") or 0)))
    x1 = min(float(item["x"]) for item in ordered)
    y1 = min(float(item["y"]) for item in ordered)
    x2 = max(float(item["x"]) + float(item["width"]) for item in ordered)
    y2 = max(float(item["y"]) + float(item["height"]) for item in ordered)
    merged = dict(ordered[0])
    font_sizes = sorted(float(item.get("font_size_px") or 0) for item in ordered if float(item.get("font_size_px") or 0) > 0)
    merged.update(
        {
            "text": _flow_join([item.get("text") for item in ordered]),
            "x": round(x1, 2),
            "y": round(y1, 2),
            "width": round(x2 - x1, 2),
            "height": round(y2 - y1, 2),
            "confidence": round(sum(float(item.get("confidence") or 0) for item in ordered) / len(ordered), 2),
            "font_size_px": round(font_sizes[len(font_sizes) // 2], 2) if font_sizes else ordered[0].get("font_size_px"),
            "source": "paragraph_group",
            "textSource": "ocr_paragraph_group",
            "lineBreakSource": "reflow_word_wrap",
            "word_wrap": True,
            "paragraphGroup": {
                "type": "same_container_multiline",
                "grouping": "column_aware_open_group",
                "lineBreakSource": "reflow_word_wrap",
                "lineCount": len(ordered),
                "lines": [paragraph_line_ref(item) for item in ordered],
            },
        }
    )
    return merged


def mark_table_rows(page: dict[str, Any]) -> None:
    """Flag glossary/table rows so they are not vertically paragraph-merged.

    A vocab card lays out "English | 中文" on the SAME row (two columns). Those
    items have a same-row partner in a different horizontal column. Stacked
    sentence lines (bubbles/cards) instead share a column at different y, so they
    are NOT flagged and still merge into flowing paragraphs.
    """
    items = [it for it in page.get("texts", []) if str(it.get("text") or "").strip()]
    page_w = int(page.get("width") or 0)
    x_tol = max(40.0, page_w * 0.02)

    def visible_len(text: str) -> int:
        return len(re.findall(r"[A-Za-z0-9㐀-鿿]", text))

    def column_size(x_left: float) -> int:
        # how many items share this left-edge column (a repeating list column)
        return sum(1 for it in items if abs(float(it.get("x") or 0) - x_left) <= x_tol)

    for a in items:
        a_text = str(a.get("text") or "")
        al = float(a.get("x") or 0)
        # A glossary entry is SHORT and sits in a column that repeats >=3 times.
        # The length gate is what separates a vocab column from a speech bubble
        # whose wrapped sentence lines also share a left edge.
        if visible_len(a_text) > 10 or column_size(al) < 3:
            continue
        a_cjk = contains_cjk(a_text)
        ah = float(a.get("height") or 0)
        acy = float(a.get("y") or 0) + ah / 2
        ar = al + float(a.get("width") or 0)
        for b in items:
            if b is a:
                continue
            b_text = str(b.get("text") or "")
            # the row partner must be a SHORT translation in the OTHER script
            if visible_len(b_text) > 10 or contains_cjk(b_text) == a_cjk:
                continue
            bh = float(b.get("height") or 0)
            bcy = float(b.get("y") or 0) + bh / 2
            if abs(acy - bcy) > max(8.0, min(ah, bh) * 0.5):
                continue  # not on the same row
            bl = float(b.get("x") or 0)
            br = bl + float(b.get("width") or 0)
            overlap = max(0.0, min(ar, br) - max(al, bl))
            if overlap < min(float(a.get("width") or 1.0), float(b.get("width") or 1.0)) * 0.3:
                a["tableRow"] = True
                break


def merge_paragraph_text_items(page: dict[str, Any]) -> None:
    page_w = int(page.get("width") or 0)
    page_h = int(page.get("height") or 0)
    items = [dict(item) for item in sorted(page.get("texts", []), key=lambda item: (float(item.get("y") or 0), float(item.get("x") or 0)))]
    candidate_indexes = [idx for idx, item in enumerate(items) if paragraph_merge_candidate(item, page_w, page_h)]
    active_groups: list[list[int]] = []

    for idx in candidate_indexes:
        current = items[idx]
        best_group_idx: int | None = None
        best_score = 10**9
        for group_idx, group in enumerate(active_groups):
            last = items[group[-1]]
            if not can_merge_paragraph_line(last, current, page_w, page_h):
                continue
            x_delta = abs(float(current.get("x") or 0) - float(last.get("x") or 0))
            y_gap = float(current.get("y") or 0) - (float(last.get("y") or 0) + float(last.get("height") or 0))
            score = x_delta + max(0.0, y_gap) * 0.6
            if score < best_score:
                best_score = score
                best_group_idx = group_idx
        if best_group_idx is None:
            active_groups.append([idx])
        else:
            active_groups[best_group_idx].append(idx)

    replacement_by_first_idx: dict[int, dict[str, Any]] = {}
    consumed: set[int] = set()
    paragraph_repairs: list[dict[str, Any]] = []
    for group in active_groups:
        if len(group) < 2:
            continue
        group_lines = [items[idx] for idx in group]
        merged = merge_paragraph_group(group_lines)
        replacement_by_first_idx[group[0]] = merged
        consumed.update(group[1:])
        paragraph_repairs.append(
            {
                "type": "paragraph_group_merge",
                "grouping": "column_aware_open_group",
                "lineCount": len(group),
                "texts": [str(item.get("text") or "") for item in group_lines],
                "reason": "Adjacent OCR rows in the same visual region, column, and style are rebuilt as one editable paragraph while preserving visible row breaks.",
            }
        )

    rebuilt: list[dict[str, Any]] = []
    for idx, item in enumerate(items):
        if idx in consumed:
            continue
        rebuilt.append(replacement_by_first_idx.get(idx, item))
    page["texts"] = sorted(rebuilt, key=lambda item: (float(item.get("y") or 0), float(item.get("x") or 0)))
    if paragraph_repairs:
        page.setdefault("layout_repairs", []).extend(paragraph_repairs)


def render_pdf(pdf: Path, out_dir: Path, dpi: int, pages_arg: str | None) -> list[dict[str, Any]]:
    out_dir.mkdir(parents=True, exist_ok=True)
    total = None
    if not pages_arg:
        try:
            info = subprocess.run(["pdfinfo", str(pdf)], check=True, capture_output=True, text=True, timeout=20)
            match = re.search(r"^Pages:\s+(\d+)", info.stdout, re.MULTILINE)
            if match:
                total = int(match.group(1))
        except (OSError, subprocess.SubprocessError):
            total = None
    pages = parse_pages(pages_arg, total)
    rendered: list[dict[str, Any]] = []
    for idx in pages:
        path = out_dir / f"slide_{idx + 1:03d}.png"
        prefix = out_dir / f"slide_{idx + 1:03d}"
        # Idempotent cache: a rendered page only changes if the source PDF/page
        # or DPI changes. Skip the (slow) pdftoppm call when the PNG already
        # exists so re-running layout/text does not re-render every page.
        if path.exists() and path.stat().st_size > 0:
            print(f"[info] reuse cached render for page {idx + 1}", file=sys.stderr)
            rendered.append({"page_number": idx + 1, "image": str(path)})
            continue
        print(f"[info] render page {idx + 1}/{len(pages)}", file=sys.stderr)
        subprocess.run(
            [
                "pdftoppm",
                "-png",
                "-singlefile",
                "-r",
                str(dpi),
                "-f",
                str(idx + 1),
                "-l",
                str(idx + 1),
                str(pdf),
                str(prefix),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )
        generated = prefix.with_suffix(".png")
        if generated != path and generated.exists():
            generated.replace(path)
        rendered.append({"page_number": idx + 1, "image": str(path)})
    return rendered


_THIS_DIR = str(Path(__file__).resolve().parent)
if _THIS_DIR not in sys.path:
    sys.path.insert(0, _THIS_DIR)
from style_probe import analyze_text_region, text_color_from_region  # noqa: E402,F401


def style_evidence_from_region(image: Image.Image, box: tuple[int, int, int, int]) -> dict[str, Any]:
    # Polarity-agnostic ink evidence via the shared style probe. The old version
    # treated `luminance < 150` as ink, which measured the background (not the
    # text) on light-on-dark regions and corrupted width-fit and bold decisions.
    probe = analyze_text_region(image, box)
    return {
        "inkDensity": round(float(probe.get("ink_density") or 0.0), 5),
        "tightInkDensity": round(float(probe.get("tight_ink_density") or 0.0), 5),
        "tightInkWidth": round(float(probe.get("glyph_width") or 0.0), 2),
        "tightInkHeight": round(float(probe.get("glyph_height") or 0.0), 2),
    }


def apply_visual_style_evidence(page: dict[str, Any]) -> None:
    from PIL import Image

    image = Image.open(page["image"]).convert("RGB")
    page_w = int(page.get("width") or image.width)
    page_h = int(page.get("height") or image.height)
    for item in page.get("texts", []):
        try:
            x1 = max(0, int(round(float(item["x"]))))
            y1 = max(0, int(round(float(item["y"]))))
            x2 = min(page_w, int(round(float(item["x"]) + float(item["width"]))))
            y2 = min(page_h, int(round(float(item["y"]) + float(item["height"]))))
        except (KeyError, TypeError, ValueError):
            continue
        if x2 <= x1 or y2 <= y1:
            continue
        evidence = style_evidence_from_region(image, (x1, y1, x2, y2))
        item["styleEvidence"] = evidence
        role = str(item.get("role") or classify_text_role(item, page_w, page_h))
        height = float(item.get("height") or 0)
        tight_density = float(evidence.get("tightInkDensity") or 0)
        bold = role == "title" or (height >= page_h * 0.052 and tight_density >= 0.27)
        item["font_bold"] = bool(bold)
        item["fontWeightSource"] = "visual_ink_density" if bold else "default_regular"
        item["fontWeightConfidence"] = round(min(1.0, tight_density / 0.36), 3)


def list_like_text(text: str) -> bool:
    clean = text.strip()
    return bool(clean.startswith(("•", "-", "*")) or ":" in clean or "：" in clean)


def apply_group_style_consistency(page: dict[str, Any]) -> None:
    page_w = int(page.get("width") or 0)
    page_h = int(page.get("height") or 0)
    candidates: list[dict[str, Any]] = []
    for item in page.get("texts", []):
        text = str(item.get("text") or "")
        y = float(item.get("y") or 0)
        height = float(item.get("height") or 0)
        if y > page_h * 0.72:
            continue
        if height < page_h * 0.045:
            continue
        if not list_like_text(text):
            continue
        candidates.append(item)

    groups: list[list[dict[str, Any]]] = []
    for item in sorted(candidates, key=lambda value: (float(value.get("x") or 0), float(value.get("y") or 0))):
        placed = False
        x = float(item.get("x") or 0)
        for group in groups:
            gx = sorted(float(value.get("x") or 0) for value in group)[len(group) // 2]
            if abs(x - gx) <= max(70.0, page_w * 0.04):
                group.append(item)
                placed = True
                break
        if not placed:
            groups.append([item])

    for group_idx, group in enumerate(groups, start=1):
        if len(group) < 3:
            continue
        densities = sorted(float((item.get("styleEvidence") or {}).get("tightInkDensity") or 0) for item in group)
        median_density = densities[len(densities) // 2]
        group_bold = median_density >= 0.24 or sum(1 for item in group if item.get("font_bold")) >= len(group) / 2
        group_id = f"list_like_column_{group_idx}"
        for item in group:
            item["font_bold"] = bool(group_bold)
            item["fontWeightSource"] = "group_style_consistency"
            item["fontWeightConfidence"] = round(min(1.0, median_density / 0.34), 3)
            item["styleGroup"] = {
                "id": group_id,
                "type": "list_like_column",
                "size": len(group),
                "medianTightInkDensity": round(median_density, 5),
            }


def style_groups(page: dict[str, Any]) -> list[list[dict[str, Any]]]:
    groups_by_key: dict[str, list[dict[str, Any]]] = {}
    for item in page.get("texts", []):
        group = item.get("styleGroup") or {}
        group_id = group.get("id")
        if not group_id:
            continue
        groups_by_key.setdefault(str(group_id), []).append(item)
    return [group for group in groups_by_key.values() if len(group) >= 3]


def apply_group_typography_consistency(page: dict[str, Any]) -> None:
    for group in style_groups(page):
        font_sizes = sorted(float(item.get("font_size_px") or 0) for item in group if float(item.get("font_size_px") or 0) > 0)
        if not font_sizes:
            continue
        median_size = font_sizes[len(font_sizes) // 2]
        for item in group:
            item["font_size_px"] = round(median_size, 2)
            item["fontSizeSource"] = "group_typography_consistency"
            item["fontSizeLocked"] = True
            item.setdefault("typographyGroup", {})
            item["typographyGroup"].update(
                {
                    "type": (item.get("styleGroup") or {}).get("type"),
                    "size": len(group),
                    "medianFontSizePx": round(median_size, 2),
                }
            )


def apply_textbox_metrics(page: dict[str, Any]) -> None:
    for item in page.get("texts", []):
        role = str(item.get("role") or "")
        text = str(item.get("text") or "")
        if item.get("styleGroup"):
            scale = 1.02
            line_spacing = 0.95
        elif role == "title":
            scale = 1.06
            line_spacing = 0.95
        elif len(text) > 42:
            scale = 1.08
            line_spacing = 0.95
        else:
            scale = 1.04
            line_spacing = 0.98
        item["textBoxHeightScale"] = scale
        item["lineSpacing"] = line_spacing
        item["textboxMetricsSource"] = "role_and_group_typography"


def fit_text_size_to_box(page: dict[str, Any]) -> None:
    """Shrink any text whose wrapped lines would overflow its box.

    OCR boxes are tight to a single ink line, so a font sized purely from glyph
    height can wrap (narrow box) or exceed the box height, causing the overflow
    and overlap seen in tight bilingual cards/bubbles. This is a pure safety cap:
    pick the largest size (<= current) at which the text, wrapped to the box
    width, still fits the box height. Items that already fit are left untouched.
    """
    import math

    for item in page.get("texts", []):
        text = str(item.get("text") or "")
        if not text.strip():
            continue
        box_w = float(item.get("width") or 0)
        scale = float(item.get("textBoxHeightScale") or 1.04)
        box_h = float(item.get("height") or 0) * scale
        cur = float(item.get("font_size_px") or 0)
        if box_w <= 0 or box_h <= 0 or cur <= 0:
            continue
        family = item.get("font_family") or font_for_text(text)
        bold = bool(item.get("font_bold"))
        # px line box per text line ~= em * 1.25 (leading); be slightly safe.
        line_box = max(1.05, float(item.get("lineSpacing") or 1.0) * 1.22)
        hard_lines = text.split("\n") or [text]
        chosen = cur
        for trial in (cur * s for s in (1.0, 0.94, 0.88, 0.82, 0.76, 0.70, 0.64, 0.58, 0.52, 0.46, 0.40)):
            total_lines = 0
            for ln in hard_lines:
                if not ln.strip():
                    total_lines += 1
                    continue
                metrics = render_font_metrics(ln, family, trial, bold)
                width = metrics["width"] if metrics else len(ln) * trial * 0.6
                total_lines += max(1, math.ceil(width / max(1.0, box_w)))
            if total_lines * trial * line_box <= box_h * 1.02:
                chosen = trial
                break
        if chosen < cur - 0.5:
            item["font_size_px"] = round(chosen, 2)
            item["fontSizeSource"] = "fit_to_box_cap"
            item["fontSizePreCap"] = round(cur, 2)


def font_file_for_family(font_family: str, bold: bool) -> str | None:
    files = FONT_FILES.get(font_family) or ()
    if not files:
        return None
    expanded = [Path(path).expanduser() for path in files]
    if bold and len(expanded) > 1 and expanded[1].exists():
        return str(expanded[1])
    if expanded[0].exists():
        return str(expanded[0])
    return None


def render_font_metrics(text: str, font_family: str, font_size_px: float, bold: bool) -> dict[str, float] | None:
    from PIL import Image, ImageDraw, ImageFont
    import numpy as np

    font_path = font_file_for_family(font_family, bold)
    if not font_path:
        return None
    try:
        font = ImageFont.truetype(font_path, max(6, int(round(font_size_px))))
    except OSError:
        return None
    bbox = ImageDraw.Draw(Image.new("L", (16, 16))).textbbox((0, 0), text, font=font)
    width = max(1, bbox[2] - bbox[0])
    height = max(1, bbox[3] - bbox[1])
    canvas = Image.new("L", (width + 12, height + 12), 255)
    draw = ImageDraw.Draw(canvas)
    draw.text((6 - bbox[0], 6 - bbox[1]), text, font=font, fill=0)
    arr = np.asarray(canvas, dtype=np.uint8)
    ink = arr < 210
    if not ink.any():
        return {"width": float(width), "height": float(height), "tightInkWidth": 0.0, "tightInkHeight": 0.0, "tightInkDensity": 0.0}
    ys, xs = np.where(ink)
    tight = ink[ys.min() : ys.max() + 1, xs.min() : xs.max() + 1]
    return {
        "width": float(xs.max() - xs.min() + 1),
        "height": float(ys.max() - ys.min() + 1),
        "tightInkWidth": float(xs.max() - xs.min() + 1),
        "tightInkHeight": float(ys.max() - ys.min() + 1),
        "tightInkDensity": float(tight.mean()),
    }


def font_fit_candidates(item: dict[str, Any]) -> tuple[str, ...]:
    text = str(item.get("text") or "")
    role = str(item.get("role") or "")
    if contains_cjk(text):
        current = str(item.get("font_family") or DEFAULT_CJK_FONT)
        return (current if current in FONT_FILES else DEFAULT_CJK_FONT,)
    if role == "title":
        return ("Comic Sans MS", "Chalkboard SE", "Arial", "Times New Roman")
    current = str(item.get("font_family") or font_for_text(text))
    candidates = [current if current in FONT_FILES else FALLBACK_LATIN_FONT, FALLBACK_LATIN_FONT]
    return tuple(dict.fromkeys(candidates))


def fit_font_for_item(item: dict[str, Any]) -> dict[str, Any] | None:
    text = str(item.get("text") or "").strip()
    if not text:
        return None
    fit_text = max([line.strip() for line in text.splitlines() if line.strip()] or [text], key=len)
    if contains_cjk(text) and bool(item.get("styleGroup")):
        return None
    evidence = item.get("styleEvidence") or {}
    target_width = float(evidence.get("tightInkWidth") or item.get("width") or 0)
    base_size = float(item.get("font_size_px") or 0)
    if target_width <= 0 or base_size <= 0:
        return None
    target_density = float(evidence.get("tightInkDensity") or 0.24)
    original_bold = bool(item.get("font_bold"))
    can_fit_weight = (str(item.get("role") or "") == "title" and not contains_cjk(text)) or (
        bool(item.get("styleGroup")) and not contains_cjk(text)
    )
    best: dict[str, Any] | None = None
    role = str(item.get("role") or "")
    for family in font_fit_candidates(item):
        trial_bold_values = (False, True) if can_fit_weight else (original_bold,)
        for trial_bold in trial_bold_values:
            if role == "title":
                scales = (0.82, 0.88, 0.94, 1.0, 1.06, 1.12)
            elif contains_cjk(text):
                scales = (0.84, 0.90, 0.96, 1.0)
            else:
                scales = (0.72, 0.78, 0.84, 0.90, 0.96, 1.0)
            for scale in scales:
                trial_size = max(6.0, base_size * scale)
                metrics = render_font_metrics(fit_text, family, trial_size, trial_bold)
                if not metrics:
                    continue
                width_ratio = metrics["width"] / max(1.0, target_width)
                width_error = abs(width_ratio - 1.0)
                density_error = abs(float(metrics["tightInkDensity"]) - target_density)
                penalty = 0.0
                if family == "Times New Roman" and (contains_cjk(text) or str(item.get("role") or "") != "title"):
                    penalty += 0.5
                if family in {"Comic Sans MS", "Chalkboard SE"} and str(item.get("role") or "") != "title":
                    penalty += 0.08
                if trial_bold and target_density < 0.19:
                    penalty += 0.06
                score = width_error * 0.62 + density_error * 0.38 + penalty
                candidate = {
                    "fontFamily": family,
                    "fontSizePx": round(trial_size, 2),
                    "fontBold": bool(trial_bold),
                    "score": round(score, 5),
                    "targetWidth": round(target_width, 2),
                    "renderWidth": round(float(metrics["width"]), 2),
                    "widthRatio": round(width_ratio, 5),
                    "renderTightInkDensity": round(float(metrics["tightInkDensity"]), 5),
                    "targetTightInkDensity": round(target_density, 5),
                }
                if best is None or candidate["score"] < best["score"]:
                    best = candidate
    if best and role != "title":
        if float(best.get("widthRatio") or 0) > 1.35 or float(best.get("score") or 0) > 0.32:
            return None
    return best


def apply_font_fit(page: dict[str, Any]) -> None:
    for item in page.get("texts", []):
        fit = fit_font_for_item(item)
        if not fit:
            continue
        old_family = item.get("font_family")
        old_size = item.get("font_size_px")
        old_bold = item.get("font_bold")
        item["font_family"] = fit["fontFamily"]
        item["font_size_px"] = fit["fontSizePx"]
        item["font_bold"] = bool(fit["fontBold"])
        item["fontFit"] = {
            **fit,
            "previousFontFamily": old_family,
            "previousFontSizePx": old_size,
            "previousFontBold": old_bold,
            "source": "rendered_width_and_ink_density",
        }
        item["fontFamilySource"] = "render_fit_approved_pool"
        if item.get("fontSizeSource") != "group_typography_consistency":
            item["fontSizeSource"] = "render_fit_compensation"

    for group in style_groups(page):
        families: dict[str, int] = {}
        bold_values: dict[bool, int] = {}
        sizes: list[float] = []
        for item in group:
            families[str(item.get("font_family") or "")] = families.get(str(item.get("font_family") or ""), 0) + 1
            bold_values[bool(item.get("font_bold"))] = bold_values.get(bool(item.get("font_bold")), 0) + 1
            sizes.append(float(item.get("font_size_px") or 0))
        if not families or not sizes:
            continue
        family = max(families.items(), key=lambda entry: entry[1])[0]
        group_bold = max(bold_values.items(), key=lambda entry: entry[1])[0]
        median_size = sorted(sizes)[len(sizes) // 2]
        for item in group:
            item["font_family"] = family
            item["font_size_px"] = round(median_size, 2)
            item["font_bold"] = bool(group_bold)
            item["fontFamilySource"] = "group_render_fit_consistency"
            item["fontSizeSource"] = "group_render_fit_consistency"
            item["fontWeightSource"] = "group_render_fit_consistency"
            item.setdefault("typographyGroup", {})
            item["typographyGroup"].update(
                {
                    "fontFamily": family,
                    "fontBold": bool(group_bold),
                    "renderFitMedianFontSizePx": round(median_size, 2),
                }
            )


def sample_background_color(image: Image.Image, box: tuple[int, int, int, int], pad: int = 12) -> tuple[int, int, int]:
    x1, y1, x2, y2 = box
    w, h = image.size
    regions = [
        (max(0, x1 - pad), max(0, y1 - pad), min(w, x2 + pad), max(0, y1)),
        (max(0, x1 - pad), min(h, y2), min(w, x2 + pad), min(h, y2 + pad)),
        (max(0, x1 - pad), max(0, y1), max(0, x1), min(h, y2)),
        (min(w, x2), max(0, y1), min(w, x2 + pad), min(h, y2)),
    ]
    samples: list[tuple[int, int, int]] = []
    for region in regions:
        if region[2] <= region[0] or region[3] <= region[1]:
            continue
        samples.extend(list(image.crop(region).convert("RGB").getdata()))
    if not samples:
        return (255, 255, 255)
    bright = [p for p in samples if sum(p) > 540]
    pool = bright or samples
    return tuple(int(sum(p[i] for p in pool) / len(pool)) for i in range(3))


def default_paddle_python() -> Path | None:
    root = Path(__file__).resolve().parents[1]
    candidates = [
        os.environ.get("PADDLEOCR_PYTHON"),
        str(root / ".venv" / "bin" / "python"),
        str(root / ".venv-paddleocr" / "bin" / "python"),
        str(Path.cwd() / ".venv" / "bin" / "python"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return Path(candidate)
    return None


def run_paddle_worker_batch(python: Path, worker: Path, rendered: list[dict[str, Any]], timeout: int) -> list[dict[str, Any]]:
    if not rendered:
        return []
    payload = json.dumps({"images": rendered}, ensure_ascii=False)
    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    env["PYTHONUNBUFFERED"] = "1"
    try:
        result = subprocess.run(
            [str(python), str(worker)],
            input=payload,
            check=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            env=env,
        )
    except subprocess.CalledProcessError as exc:
        stderr = (exc.stderr or "").strip()
        stdout = (exc.stdout or "").strip()
        raise RuntimeError(
            "PaddleOCR worker failed "
            f"returncode={exc.returncode} stderr={stderr[-4000:]} stdout={stdout[-1000:]}"
        ) from exc
    except subprocess.TimeoutExpired as exc:
        stderr = (exc.stderr or b"")
        if isinstance(stderr, bytes):
            stderr = stderr.decode("utf-8", errors="replace")
        raise RuntimeError(f"PaddleOCR worker timed out after {timeout}s stderr={str(stderr)[-4000:]}") from exc
    json_line = ""
    for line in reversed(result.stdout.splitlines()):
        if line.strip().startswith("{"):
            json_line = line.strip()
            break
    if not json_line:
        raise RuntimeError(f"PaddleOCR worker produced no JSON. stderr={result.stderr.strip()}")
    data = json.loads(json_line)
    if not data.get("ok"):
        raise RuntimeError("PaddleOCR worker returned ok=false")
    return data.get("slides") or []


def run_paddle_worker(rendered: list[dict[str, Any]], timeout: int, batch_size: int) -> list[dict[str, Any]]:
    python = default_paddle_python()
    if not python:
        raise RuntimeError("PaddleOCR Python runtime not found. Set PADDLEOCR_PYTHON.")
    worker = Path(__file__).with_name("ocr_paddle_worker.py")
    if batch_size <= 0:
        batch_size = max(1, len(rendered))
    else:
        batch_size = max(1, int(batch_size))
    slides: list[dict[str, Any]] = []
    total_batches = (len(rendered) + batch_size - 1) // batch_size
    for batch_index, start in enumerate(range(0, len(rendered), batch_size), start=1):
        chunk = rendered[start : start + batch_size]
        first_page = chunk[0]["page_number"]
        last_page = chunk[-1]["page_number"]
        print(
            f"[info] OCR batch {batch_index}/{total_batches} pages {first_page}-{last_page} with {python}",
            file=sys.stderr,
        )
        slides.extend(run_paddle_worker_batch(python, worker, chunk, timeout))
    return slides


def remove_notebooklm_watermark(
    image,
    x_frac: float = 0.82,
    y_frac: float = 0.965,
    diff_thresh: int = 26,
):
    """Erase the NotebookLM export watermark in the bottom-right corner.

    PaddleOCR does not pick up the light-grey "NotebookLM" mark, so it is located
    visually: only the very bottom strip (below y_frac) is searched, which sits
    BELOW any QR code in that corner, so the QR is never touched. The watermark
    ink bbox is filled with the locally-sampled background colour so it blends in.
    """
    import numpy as np
    from PIL import Image as _Image

    W, H = image.size
    rx, ry = int(W * x_frac), int(H * y_frac)
    if rx >= W - 4 or ry >= H - 4:
        return image
    region = np.asarray(image.crop((rx, ry, W, H)).convert("RGB"), dtype=np.int16)
    rh, rw = region.shape[:2]
    # background colour = median of the strip's left margin (clean page edge)
    margin = region[:, : max(3, rw // 8)].reshape(-1, 3)
    bg = np.median(margin, axis=0)
    diff = np.abs(region - bg).sum(axis=2)
    ink = diff > diff_thresh
    ys, xs = np.where(ink)
    if ys.size < 8:
        return image  # no watermark ink found -> leave the page untouched
    pad = 8
    bx1, by1 = max(0, int(xs.min()) - pad), max(0, int(ys.min()) - pad)
    bx2, by2 = min(rw, int(xs.max()) + pad), min(rh, int(ys.max()) + pad)
    patch = _Image.new("RGB", (bx2 - bx1, by2 - by1), tuple(int(c) for c in bg))
    image.paste(patch, (rx + bx1, ry + by1))
    return image


def strip_watermark_from_background(page: dict[str, Any]) -> None:
    """Apply watermark removal to this page's generated background image in place."""
    from PIL import Image

    bg_path = page.get("clean_background")
    if not bg_path or not Path(bg_path).exists():
        return
    if str(bg_path) == str(page.get("image")):
        return  # never edit the shared raw render (model-clean -> original fallback)
    try:
        img = Image.open(bg_path).convert("RGB")
        remove_notebooklm_watermark(img)
        img.save(bg_path)
        page["watermarkStripped"] = True
    except OSError:
        return


def clean_background(page: dict[str, Any], out_path: Path, expand_px: int) -> None:
    from PIL import Image, ImageFilter

    image = Image.open(page["image"]).convert("RGB")
    # Final editable text boxes. local-clean must only erase regions that will be
    # redrawn as editable text; a masked region with NO surviving editable node
    # (e.g. a big decorative number OCR'd into mask_words but later dropped) would
    # otherwise be covered AND never re-added -> the element silently disappears.
    final_boxes = []
    for t in page.get("texts", []):
        try:
            fx, fy = float(t["x"]), float(t["y"])
            final_boxes.append((fx, fy, fx + float(t["width"]), fy + float(t["height"])))
        except (KeyError, TypeError, ValueError):
            continue

    def _has_replacement(bx1: float, by1: float, bx2: float, by2: float) -> bool:
        area = max(1.0, (bx2 - bx1) * (by2 - by1))
        for fx1, fy1, fx2, fy2 in final_boxes:
            ox = max(0.0, min(bx2, fx2) - max(bx1, fx1))
            oy = max(0.0, min(by2, fy2) - max(by1, fy1))
            if ox * oy / area >= 0.3:
                return True
        return False

    for item in page.get("mask_words") or page.get("texts") or []:
        role = str(item.get("role") or "")
        item_expand = max(2, int(expand_px * (0.55 if role == "title" else 1.0)))
        x1 = max(0, int(item["x"]) - item_expand)
        y1 = max(0, int(item["y"]) - item_expand)
        x2 = min(image.width, int(item["x"] + item["width"]) + item_expand)
        y2 = min(image.height, int(item["y"] + item["height"]) + item_expand)
        if x2 <= x1 or y2 <= y1:
            continue
        if not _has_replacement(x1, y1, x2, y2):
            continue  # nothing will be redrawn here; keep the original pixels
        fill = sample_background_color(image, (x1, y1, x2, y2))
        patch = Image.new("RGB", (x2 - x1, y2 - y1), fill)
        mask = Image.new("L", (x2 - x1, y2 - y1), 255)
        if min(mask.size) > 6:
            mask = mask.filter(ImageFilter.GaussianBlur(radius=max(1, item_expand // 2)))
        image.paste(patch, (x1, y1), mask)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(out_path)
    page["clean_background"] = str(out_path)


def model_clean_prompt(page: dict[str, Any]) -> str:
    texts = [str(item.get("text") or "").strip() for item in page.get("texts", []) if str(item.get("text") or "").strip()]
    text_list = "\n".join(f"- {text}" for text in texts[:30])
    return (
        "Image editing task. Use the input image as the exact base image and keep the original geometry locked. "
        "Preserve canvas size, aspect ratio, composition, object positions, visual hierarchy, containers, panels, bubbles, cards, icons, illustrations, charts, decorations, colors, shadows, and background texture. "
        "Do not redesign, regenerate, beautify, crop, zoom, shift, move, replace, or reinterpret any non-text visual element. "
        "Only remove readable text pixels and perform the smallest necessary local background repair where those text pixels were. "
        "Keep text containers in their original shape, edge, shadow, size, and position. "
        "Do not add replacement text or new visual elements. "
        "If a listed text is inside a container, make only that original text area clean; do not redraw the container.\n\n"
        f"Readable text to remove:\n{text_list}"
    )


def normalize_model_clean_canvas(original_image: Path, cleaned_image: Path) -> tuple[Path, dict[str, Any]]:
    from PIL import Image

    original = Image.open(original_image)
    cleaned = Image.open(cleaned_image)
    if original.size == cleaned.size:
        return cleaned_image, {
            "sizeMatch": True,
            "inputSize": [original.size[0], original.size[1]],
            "outputSize": [cleaned.size[0], cleaned.size[1]],
            "normalizedToInputSize": False,
            "note": "Size match passed. Manual or visual-diff QA is still required to detect object/container drift.",
        }

    original_ratio = original.size[0] / original.size[1]
    cleaned_ratio = cleaned.size[0] / cleaned.size[1]
    ratio_delta = abs(original_ratio - cleaned_ratio) / original_ratio
    if ratio_delta > 0.01:
        raise RuntimeError(
            "model-clean geometry failed: output aspect ratio "
            f"{cleaned.size[0]}x{cleaned.size[1]} differs from input {original.size[0]}x{original.size[1]}"
        )

    normalized = cleaned_image.with_name(f"{cleaned_image.stem}.normalized{cleaned_image.suffix}")
    cleaned.convert("RGB").resize(original.size, Image.Resampling.LANCZOS).save(normalized)
    return normalized, {
        "sizeMatch": False,
        "inputSize": [original.size[0], original.size[1]],
        "outputSize": [cleaned.size[0], cleaned.size[1]],
        "normalizedToInputSize": True,
        "normalizedImage": str(normalized),
        "aspectRatioDelta": ratio_delta,
        "note": "Output kept the same aspect ratio but changed pixel size, so it was normalized back to the input canvas. Manual or visual-diff QA is still required to detect object/container drift.",
    }


def text_mask_boxes(page: dict[str, Any], expand_px: int = 10) -> list[tuple[int, int, int, int]]:
    width = int(page.get("width") or 0)
    height = int(page.get("height") or 0)
    boxes: list[tuple[int, int, int, int]] = []
    for item in page.get("mask_words") or page.get("texts") or []:
        try:
            x1 = max(0, int(round(float(item["x"]))) - expand_px)
            y1 = max(0, int(round(float(item["y"]))) - expand_px)
            x2 = min(width, int(round(float(item["x"]) + float(item["width"]))) + expand_px)
            y2 = min(height, int(round(float(item["y"]) + float(item["height"]))) + expand_px)
        except (KeyError, TypeError, ValueError):
            continue
        if x2 > x1 and y2 > y1:
            boxes.append((x1, y1, x2, y2))
    return boxes


def visual_diff_qa(original_image: Path, cleaned_image: Path, page: dict[str, Any]) -> dict[str, Any]:
    """Measure model-clean drift outside OCR text regions.

    This is intentionally a conservative QA signal, not a hard visual judge.
    Text regions are expected to change; non-text regions should stay stable.
    """

    from PIL import Image, ImageDraw
    import numpy as np

    original = Image.open(original_image).convert("RGB")
    cleaned = Image.open(cleaned_image).convert("RGB")
    if original.size != cleaned.size:
        return {
            "status": "fail",
            "reason": "size_mismatch_after_normalization",
            "inputSize": [original.size[0], original.size[1]],
            "outputSize": [cleaned.size[0], cleaned.size[1]],
        }

    mask = Image.new("L", original.size, 0)
    draw = ImageDraw.Draw(mask)
    for box in text_mask_boxes(page):
        draw.rectangle(box, fill=255)

    original_arr = np.asarray(original, dtype=np.int16)
    cleaned_arr = np.asarray(cleaned, dtype=np.int16)
    mask_arr = np.asarray(mask, dtype=np.uint8) > 0
    protected = ~mask_arr
    protected_count = int(protected.sum())
    if protected_count == 0:
        return {"status": "review", "reason": "no_non_text_pixels_for_comparison"}

    diff = np.abs(original_arr - cleaned_arr).mean(axis=2)
    protected_diff = diff[protected]
    mean_abs_diff = float(protected_diff.mean())
    p95_abs_diff = float(np.percentile(protected_diff, 95))
    changed_ratio = float((protected_diff > 35).sum() / protected_count)
    text_mask_ratio = float(mask_arr.sum() / mask_arr.size)

    status = "pass"
    reason = "non_text_regions_stable"
    if changed_ratio > 0.18 or mean_abs_diff > 22:
        status = "review"
        reason = "non_text_regions_changed_noticeably"
    if changed_ratio > 0.35 or mean_abs_diff > 38:
        status = "fail"
        reason = "non_text_regions_changed_too_much"

    return {
        "status": status,
        "reason": reason,
        "meanAbsDiff": round(mean_abs_diff, 3),
        "p95AbsDiff": round(p95_abs_diff, 3),
        "changedPixelRatio": round(changed_ratio, 5),
        "textMaskRatio": round(text_mask_ratio, 5),
        "changedPixelThreshold": 35,
        "protectedPixelCount": protected_count,
        "note": "Text-mask regions are ignored. This QA checks whether the image model drifted non-text content such as containers, icons, illustrations, cards, or composition.",
    }


def text_background_is_uniform(
    page: dict[str, Any],
    std_thresh: float = 26.0,
    bad_frac: float = 0.18,
    margin: int = 10,
) -> bool:
    """Decide if local-clean is safe for this page (vs needing model-clean).

    local-clean covers each text region with sampled neighbouring colour, so it
    is invisible only when the background AROUND the text is uniform. For every
    text box we measure the luminance spread of the surrounding margin frame; if
    enough boxes sit on a textured/illustrated/dark-varied background, the page
    needs model-clean. White/flat-colour decks come back uniform -> local-clean
    (fast, no model call). Bias is safe: only flips to model-clean on clear
    evidence of texture.
    """
    import numpy as np
    from PIL import Image

    try:
        img = Image.open(page["image"]).convert("RGB")
    except (KeyError, OSError):
        return True
    lum = np.asarray(img, dtype=np.float32) @ np.array([0.299, 0.587, 0.114], dtype=np.float32)
    H, W = lum.shape
    bad = 0
    total = 0
    for it in page.get("texts", []):
        try:
            x, y = float(it["x"]), float(it["y"])
            w, h = float(it["width"]), float(it["height"])
        except (KeyError, TypeError, ValueError):
            continue
        x1, y1 = max(0, int(x - margin)), max(0, int(y - margin))
        x2, y2 = min(W, int(x + w + margin)), min(H, int(y + h + margin))
        if x2 - x1 < 6 or y2 - y1 < 6:
            continue
        crop = lum[y1:y2, x1:x2]
        frame = np.ones(crop.shape, dtype=bool)
        iy1, iy2 = max(0, int(y) - y1), min(crop.shape[0], int(y + h) - y1)
        ix1, ix2 = max(0, int(x) - x1), min(crop.shape[1], int(x + w) - x1)
        frame[iy1:iy2, ix1:ix2] = False  # exclude the glyph box; keep the surround
        bg = crop[frame]
        if bg.size < 20:
            bg = crop.reshape(-1)
        total += 1
        if float(bg.std()) > std_thresh:
            bad += 1
    if total == 0:
        return True
    return (bad / total) < bad_frac


def model_clean_background(
    page: dict[str, Any],
    out_dir: Path,
    provider: str,
    model: str,
    base_url: str,
    api_key_env: str,
    timeout: int,
    size: str,
    insecure: bool,
) -> None:
    script = Path(__file__).with_name("repair_background_with_image_model.py")
    page_out = out_dir / f"slide_{int(page['page_number']):03d}"
    # Idempotent cache: if a cleaned background for this page+model already
    # exists, reuse it and skip the (slow, billed) image-model call. Lets us
    # re-run layout/text with new code without regenerating backgrounds.
    for cand in (
        page_out / f"{model}.clean_background.normalized.png",
        page_out / f"{model}.clean_background.png",
    ):
        if cand.exists():
            page["clean_background"] = str(cand)
            page["model_clean_geometry_qa"] = {"status": "reused"}
            page["model_clean_visual_qa"] = {
                "status": "reused",
                "reason": "cached_clean_background",
            }
            print(
                f"[info] reuse cached clean background for page "
                f"{page['page_number']}: {cand.name}",
                file=sys.stderr,
            )
            return
    page_out.mkdir(parents=True, exist_ok=True)
    if provider == "codex-image":
        prompt = model_clean_prompt(page)
        request = {
            "ok": False,
            "provider": "codex-image",
            "requires_codex_image_tool": True,
            "page": page.get("page_number"),
            "image": str(page["image"]),
            "expected_output": str(page_out / f"{model}.clean_background.png"),
            "prompt": prompt,
            "instructions": [
                "Use the Codex built-in image editing tool on the source image.",
                "Edit the original image directly; do not regenerate or redesign the slide.",
                "Remove only the OCR-listed text pixels and keep containers, icons, illustrations, cards, panels, composition, aspect ratio, and canvas geometry unchanged.",
                "Save the resulting clean background to expected_output, then rerun the PPTX rebuild with this cached background.",
            ],
        }
        request_json = page_out / f"{model}.codex_image_request.json"
        request_md = page_out / f"{model}.codex_image_request.md"
        request_json.write_text(json.dumps(request, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        request_md.write_text(
            "\n".join(
                [
                    "# Codex Image Clean Request",
                    "",
                    f"- Page: {page.get('page_number')}",
                    f"- Source image: `{page['image']}`",
                    f"- Expected output: `{request['expected_output']}`",
                    "",
                    "## Prompt",
                    "",
                    prompt,
                    "",
                    "## Rules",
                    "",
                    "- Edit the original image directly.",
                    "- Remove only the OCR-listed text pixels.",
                    "- Preserve all non-text visuals and canvas geometry.",
                ]
            )
            + "\n",
            encoding="utf-8",
        )
        page["codex_image_clean_request"] = str(request_json)
        raise RuntimeError(f"codex-image clean requires Codex image tool orchestration: {request_json}")
    cmd = [
        sys.executable,
        str(script),
        "--image",
        str(page["image"]),
        "--output-dir",
        str(page_out),
        "--provider",
        provider,
        "--model",
        model,
        "--base-url",
        base_url,
        "--api-key-env",
        api_key_env,
        "--prompt",
        model_clean_prompt(page),
        "--timeout",
        str(timeout),
        "--size",
        size,
    ]
    if insecure:
        cmd.append("--insecure")
    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    if not env.get(api_key_env):
        raise RuntimeError(f"missing API key env for model-clean: {api_key_env}")
    try:
        result = subprocess.run(cmd, check=True, capture_output=True, text=True, timeout=timeout + 30, env=env)
    except subprocess.CalledProcessError as exc:
        raise RuntimeError(
            "model-clean subprocess failed "
            f"returncode={exc.returncode} stdout={exc.stdout[-1000:] if exc.stdout else ''} "
            f"stderr={exc.stderr[-1000:] if exc.stderr else ''}"
        ) from exc
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(
            "model-clean subprocess timed out "
            f"timeout={exc.timeout} stdout={exc.stdout[-1000:] if exc.stdout else ''} "
            f"stderr={exc.stderr[-1000:] if exc.stderr else ''}"
        ) from exc
    decoder = json.JSONDecoder()
    json_data = None
    text = result.stdout
    for idx, char in enumerate(text):
        if char != "{":
            continue
        try:
            candidate, end = decoder.raw_decode(text[idx:])
        except json.JSONDecodeError:
            continue
        if text[idx + end :].strip() == "":
            json_data = candidate
            break
    if json_data is None:
        raise RuntimeError(f"model-clean produced no JSON. stderr={result.stderr.strip()}")
    normalized_image, geometry_qa = normalize_model_clean_canvas(Path(page["image"]), Path(json_data["image"]))
    visual_qa = visual_diff_qa(Path(page["image"]), normalized_image, page)
    page["model_clean_geometry_qa"] = geometry_qa
    page["model_clean_visual_qa"] = visual_qa
    page["clean_background"] = str(normalized_image)
    page["model_clean_response"] = json_data.get("response")


def hex_color(color: str):
    from pptx.dml.color import RGBColor

    clean = (color or "#111111").lstrip("#")
    if len(clean) != 6:
        clean = "111111"
    return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))


def set_run_font(run, font_name: str, size_pt: float, color: str, bold: bool = False) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Pt

    run.font.name = font_name
    run.font.size = Pt(max(1, size_pt))
    run.font.color.rgb = hex_color(color)
    run.font.bold = bool(bold)
    rpr = run._r.get_or_add_rPr()
    latin = rpr.get_or_add_latin()
    latin.set("typeface", font_name)
    for tag in ("a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", font_name)


def add_text(slide, item: dict[str, Any], img_w: int, img_h: int) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    x = item["x"] / img_w * SLIDE_W_IN
    y = item["y"] / img_h * SLIDE_H_IN
    w = max(item["width"] / img_w * SLIDE_W_IN, 0.05)
    height_scale = float(item.get("textBoxHeightScale") or 1.04)
    h = max(item["height"] / img_h * SLIDE_H_IN * height_scale, 0.08)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = bool(item.get("word_wrap"))
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = float(item.get("lineSpacing") or 1.0)
    run = p.add_run()
    run.text = item["text"]
    size_pt = item["font_size_px"] / img_h * SLIDE_H_IN * 72
    font_name = item.get("font_family") or font_for_text(item["text"])
    set_run_font(run, font_name, size_pt, item.get("color", "#111111"), bool(item.get("font_bold")))


def build_pptx(layout: dict[str, Any], pptx_path: Path, background_key: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = int(SLIDE_W_IN * EMU_PER_INCH)
    prs.slide_height = int(SLIDE_H_IN * EMU_PER_INCH)
    blank = prs.slide_layouts[6]
    for page in layout["slides"]:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            page.get(background_key) or page["image"],
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        for item in page.get("texts", []):
            add_text(slide, item, page["width"], page["height"])
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)


def page_qa_summary(page: dict[str, Any]) -> dict[str, Any]:
    texts = page.get("texts", [])
    font_fit_count = sum(1 for item in texts if item.get("fontFit"))
    font_fit_families = sorted({str(item.get("font_family")) for item in texts if item.get("fontFit") and item.get("font_family")})
    paragraph_groups = [item.get("paragraphGroup") or {} for item in texts if item.get("paragraphGroup")]
    paragraph_group_count = sum(1 for group in paragraph_groups if int(group.get("lineCount") or len(group.get("lines") or [])) >= 2)
    paragraph_group_lines = sum(int(group.get("lineCount") or len(group.get("lines") or [])) for group in paragraph_groups if int(group.get("lineCount") or len(group.get("lines") or [])) >= 2)
    ocr_repairs = page.get("ocr_repairs") or []
    visual_qa = page.get("model_clean_visual_qa") or {}
    geometry_qa = page.get("model_clean_geometry_qa") or {}
    return {
        "page": page.get("page_number"),
        "textNodes": len(texts),
        "ocrRepairs": len(ocr_repairs),
        "ocrRepairTypes": sorted({str(item.get("type")) for item in ocr_repairs if item.get("type")}),
        "paragraphGroupCount": paragraph_group_count,
        "paragraphGroupLines": paragraph_group_lines,
        "fontFitCount": font_fit_count,
        "fontFitFamilies": font_fit_families,
        "backgroundVisualStatus": visual_qa.get("status"),
        "backgroundVisualReason": visual_qa.get("reason"),
        "backgroundChangedPixelRatio": visual_qa.get("changedPixelRatio"),
        "backgroundMeanAbsDiff": visual_qa.get("meanAbsDiff"),
        "backgroundGeometryNormalized": geometry_qa.get("normalizedToInputSize"),
        "backgroundGeometryAspectRatioDelta": geometry_qa.get("aspectRatioDelta"),
        "modelCleanFallback": page.get("model_clean_fallback"),
        "modelCleanError": page.get("model_clean_error"),
        "needsReview": bool(
            visual_qa.get("status") in {"review", "fail"}
            or page.get("model_clean_error")
            or len(ocr_repairs) > 0
            or any(str(item.get("text") or "").strip() == "" for item in texts)
        ),
        "sampleTexts": [str(item.get("text") or "") for item in texts[:8]],
    }


def write_qa_summary(layout: dict[str, Any], qa_path: Path) -> dict[str, Any]:
    pages = [page_qa_summary(page) for page in layout.get("slides", [])]
    payload = {
        "source": layout.get("source"),
        "engine": layout.get("engine"),
        "background": layout.get("background"),
        "pages": pages,
        "totals": {
            "slides": len(pages),
            "textNodes": sum(int(page.get("textNodes") or 0) for page in pages),
            "ocrRepairs": sum(int(page.get("ocrRepairs") or 0) for page in pages),
            "paragraphGroupCount": sum(int(page.get("paragraphGroupCount") or 0) for page in pages),
            "paragraphGroupLines": sum(int(page.get("paragraphGroupLines") or 0) for page in pages),
            "fontFitCount": sum(int(page.get("fontFitCount") or 0) for page in pages),
            "needsReview": sum(1 for page in pages if page.get("needsReview")),
            "visualPass": sum(1 for page in pages if page.get("backgroundVisualStatus") == "pass"),
            "visualReview": sum(1 for page in pages if page.get("backgroundVisualStatus") == "review"),
            "visualFail": sum(1 for page in pages if page.get("backgroundVisualStatus") == "fail"),
        },
    }
    qa_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return payload


def main() -> int:
    parser = argparse.ArgumentParser(description="Simple editable PPTX reconstruction from PDF")
    parser.add_argument("--pdf", required=True)
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--pages", help="1-based pages, e.g. 1,2,5-7")
    parser.add_argument("--dpi", type=int, default=180)
    parser.add_argument("--ocr", choices=["auto", "paddle"], default="auto")
    parser.add_argument("--lang", default="chi_sim+eng")
    parser.add_argument("--psm", type=int, default=11)
    parser.add_argument("--min-conf", type=int, default=35)
    parser.add_argument("--ocr-timeout", type=int, default=300)
    parser.add_argument("--ocr-batch-size", type=int, default=0, help="pages per OCR worker; 0 means all selected pages in one worker")
    parser.add_argument("--background", choices=["original", "clean-text", "local-clean", "model-clean", "auto"], default="auto")
    parser.add_argument("--mask-expand", type=int, default=6)
    # Default to the OpenAI image-edit endpoint for portability. Users with a
    # compatible proxy can set VISION_API_BASE_URL or pass --model-clean-base-url.
    parser.add_argument("--model-provider", choices=["codex-image", "gemini-native", "openai-image"], default="codex-image")
    parser.add_argument("--model-clean-model", default="codex-image")
    parser.add_argument("--model-clean-base-url", default=os.environ.get("VISION_API_BASE_URL", "https://api.openai.com"))
    parser.add_argument("--model-clean-api-key-env", default="VISION_API_KEY")
    parser.add_argument("--model-clean-timeout", type=int, default=240)
    parser.add_argument("--model-clean-size", default="1536x864")
    parser.add_argument("--model-clean-insecure", action="store_true")
    parser.add_argument("--model-clean-fallback", choices=["fail", "local-clean", "original"], default="local-clean")
    args = parser.parse_args()

    pdf = Path(args.pdf).expanduser().resolve()
    out_dir = Path(args.output_dir).expanduser().resolve()
    rendered_dir = out_dir / "01_rendered"
    ocr_dir = out_dir / "02_ocr"
    clean_dir = out_dir / "03_cleaned"
    pptx_dir = out_dir / "04_pptx"
    model_clean_dir = out_dir / "03_model_cleaned"

    rendered = render_pdf(pdf, rendered_dir, args.dpi, args.pages)
    print(f"[info] rendered {len(rendered)} page(s) into {rendered_dir}", file=sys.stderr)
    slides: list[dict[str, Any]] = []
    if args.ocr in {"auto", "paddle"}:
        print(f"[info] OCR {len(rendered)} page(s) with PaddleOCR worker", file=sys.stderr)
        try:
            slides = run_paddle_worker(rendered, args.ocr_timeout, args.ocr_batch_size)
        except Exception as exc:
            raise RuntimeError(f"PaddleOCR worker failed in the default flow: {exc}") from exc

    if not slides:
        raise RuntimeError("PaddleOCR produced no slides in the default flow.")

    slides = [normalize_ocr_slide(page) for page in slides]
    for page in slides:
        merge_section_numeral_prefix(page)
        merge_primary_top_band_heading(page)
        mark_table_rows(page)
        merge_cjk_short_continuations(page)
        merge_paragraph_text_items(page)
        apply_visual_style_evidence(page)
        apply_group_style_consistency(page)
        apply_group_typography_consistency(page)
        apply_font_fit(page)
        apply_textbox_metrics(page)
        fit_text_size_to_box(page)

    for parsed in slides:
        image_path = Path(parsed["image"])
        mode = args.background
        if mode == "auto":
            uniform = text_background_is_uniform(parsed)
            mode = "local-clean" if uniform else "model-clean"
            parsed["backgroundAutoChoice"] = mode
            print(
                f"[info] auto: page {parsed['page_number']} -> {mode} "
                f"({'uniform text background' if uniform else 'textured/dark text background'})",
                file=sys.stderr,
            )
        if mode in {"clean-text", "local-clean"}:
            print(f"[info] clean text background for page {parsed['page_number']}", file=sys.stderr)
            clean_background(parsed, clean_dir / image_path.name, args.mask_expand)
        elif mode == "model-clean":
            print(f"[info] model clean background for page {parsed['page_number']}", file=sys.stderr)
            try:
                model_clean_background(
                    parsed,
                    model_clean_dir,
                    args.model_provider,
                    args.model_clean_model,
                    args.model_clean_base_url,
                    args.model_clean_api_key_env,
                    args.model_clean_timeout,
                    args.model_clean_size,
                    args.model_clean_insecure,
                )
            except Exception as exc:
                parsed["model_clean_error"] = str(exc)
                parsed["model_clean_visual_qa"] = {
                    "status": "fail",
                    "reason": "model_clean_failed",
                    "error": str(exc),
                }
                print(f"[warn] model clean failed for page {parsed['page_number']}: {exc}", file=sys.stderr)
                if args.model_clean_fallback == "fail":
                    raise
                if args.model_clean_fallback == "local-clean":
                    fallback_path = clean_dir / image_path.name
                    clean_background(parsed, fallback_path, args.mask_expand)
                    parsed["model_clean_fallback"] = "local-clean"
                elif args.model_clean_fallback == "original":
                    parsed["clean_background"] = str(image_path)
                    parsed["model_clean_fallback"] = "original"
        # Strip the NotebookLM export watermark from the generated background so
        # it blends with the page. Skipped for "original" mode (no clean bg copy).
        if mode in {"clean-text", "local-clean", "model-clean"}:
            strip_watermark_from_background(parsed)

    layout = {
        "source": str(pdf),
        "engine": args.ocr,
        "dpi": args.dpi,
        "background": args.background,
        "model_clean": {
            "provider": args.model_provider,
            "model": args.model_clean_model,
            "base_url": args.model_clean_base_url,
            "api_key_env": args.model_clean_api_key_env,
        } if args.background == "model-clean" else None,
        "slides": slides,
    }
    ocr_dir.mkdir(parents=True, exist_ok=True)
    layout_path = ocr_dir / "layout.json"
    layout_path.write_text(json.dumps(layout, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    qa_path = ocr_dir / "qa_summary.json"
    qa_summary = write_qa_summary(layout, qa_path)

    pptx_path = pptx_dir / "editable_text_overlay.pptx"
    background_key = "clean_background" if args.background in {"clean-text", "local-clean", "model-clean", "auto"} else "image"
    print(f"[info] build pptx at {pptx_path}", file=sys.stderr)
    build_pptx(layout, pptx_path, background_key)

    print(
        json.dumps(
            {
                "ok": True,
                "slides": len(slides),
                "text_nodes": sum(len(page.get("texts", [])) for page in slides),
                "layout": str(layout_path),
                "qa_summary": str(qa_path),
                "pptx": str(pptx_path),
                "work_dir": str(out_dir),
                "qa_totals": qa_summary.get("totals"),
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
