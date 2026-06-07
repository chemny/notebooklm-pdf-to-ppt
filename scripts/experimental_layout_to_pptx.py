#!/usr/bin/env python3
"""Minimal PPTX renderer for experimental layout JSON.

This avoids importing editable_deck.py so experiments are not blocked by its
full OCR/font probing startup path.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5
EMU_PER_INCH = 914400
FONT_SCALE = 1.03


def contains_cjk(text: str) -> bool:
    return bool(re.search(r"[\u3400-\u9fff]", text or ""))


def hex_to_rgb(color: str | None) -> RGBColor:
    value = (color or "#111111").lstrip("#")
    if len(value) != 6:
        value = "111111"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def font_for(text: str, item: dict[str, Any]) -> str:
    family = item.get("font_family") or item.get("fontFamily")
    if family:
        return str(family)
    return "Noto Sans SC" if contains_cjk(text) else "Inter"


def align_for(value: Any) -> PP_ALIGN:
    if isinstance(value, list):
        value = value[0] if value else "LEFT"
    return {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFIED": PP_ALIGN.JUSTIFY,
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(str(value or "LEFT"), PP_ALIGN.LEFT)


def set_run_style(run, item: dict[str, Any], img_h: float) -> None:
    text = str(item.get("text") or "")
    font_name = font_for(text, item)
    run.font.name = font_name
    run.font.size = Pt(max(1.0, float(item.get("font_size_px") or 28) * SLIDE_H / img_h * 72 * FONT_SCALE))
    run.font.bold = bool(item.get("bold"))
    run.font.color.rgb = hex_to_rgb(item.get("color"))
    rpr = run._r.get_or_add_rPr()
    latin = rpr.get_or_add_latin()
    latin.set("typeface", font_name)
    for tag in ("a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", font_name)


def add_textbox(slide, item: dict[str, Any], img_w: float, img_h: float) -> None:
    x = float(item.get("x", 0)) / img_w * SLIDE_W
    y = float(item.get("y", 0)) / img_h * SLIDE_H
    w = max(float(item.get("width", 100)) / img_w * SLIDE_W, 0.15)
    h = max(float(item.get("height", 30)) / img_h * SLIDE_H, 0.12)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = str(item.get("text") or "").splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align_for(item.get("align"))
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = float(item.get("lineSpacing") or item.get("line_spacing") or 0.92)
        run = p.add_run()
        run.text = line
        set_run_style(run, item, img_h)


def build(layout: dict[str, Any], output: Path, background_key: str) -> None:
    prs = Presentation()
    prs.slide_width = int(SLIDE_W * EMU_PER_INCH)
    prs.slide_height = int(SLIDE_H * EMU_PER_INCH)
    blank = prs.slide_layouts[6]
    for page in layout.get("slides", []):
        slide = prs.slides.add_slide(blank)
        bg = page.get(background_key) or page.get("image")
        slide.shapes.add_picture(str(bg), 0, 0, width=prs.slide_width, height=prs.slide_height)
        img_w = float(page.get("width") or 3440)
        img_h = float(page.get("height") or 1920)
        for item in page.get("texts", []):
            add_textbox(slide, item, img_w, img_h)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    parser = argparse.ArgumentParser(description="Render experimental layout JSON to PPTX")
    parser.add_argument("--layout", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--background-key", default="clean_background")
    args = parser.parse_args()
    layout = json.loads(Path(args.layout).expanduser().read_text(encoding="utf-8"))
    build(layout, Path(args.output).expanduser(), args.background_key)
    print(json.dumps({"ok": True, "pptx": str(Path(args.output).expanduser().resolve())}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
