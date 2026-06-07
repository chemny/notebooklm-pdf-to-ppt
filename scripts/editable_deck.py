#!/usr/bin/env python3
"""Create a first-pass editable PPTX from image/PDF-based slide decks.

Modes:
- text-overlay: original slide image as background, OCR text boxes on top.
- clean-text: text regions are masked/fill-repaired first, then editable text
  boxes are layered on top. This is a local first pass, not generative
  inpainting.
- cover-text: original slide image as background, editable text boxes have a
  sampled opaque fill to cover the original text without damaging images.

The script accepts model-repaired layout JSON through --repair-layout. Repaired
layouts may either keep the OCR-first "texts" list or use a Codia-like
"elements" list with text/image objects.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import tempfile
import zipfile
from functools import lru_cache
from pathlib import Path
from typing import Any

import fitz
import pytesseract
from PIL import Image
from PIL import ImageDraw
from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
DEFAULT_SLIDE_W = 13.333333
DEFAULT_SLIDE_H = 7.5
PPT_RENDER_FONT_SCALE = 1.03
CJK_FONT_CANDIDATES = ("Noto Sans SC", "Source Han Sans CN", "思源黑体 CN", "Arial", "Times New Roman")
LATIN_FONT_CANDIDATES = ("Inter", "Arial", "Times New Roman")
FONT_FILE_HINTS = {
    "Noto Sans SC": ("~/Library/Fonts/NotoSansSC-Variable.ttf", "/Library/Fonts/NotoSansSC-Variable.ttf"),
    "Source Han Sans CN": (
        "~/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "~/Library/Fonts/SourceHanSansCN-Bold.otf",
        "/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "/Library/Fonts/SourceHanSansCN-Bold.otf",
    ),
    "思源黑体 CN": (
        "~/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "~/Library/Fonts/SourceHanSansCN-Bold.otf",
        "/Library/Fonts/思源黑体SourceHanSansCN-Medium.otf",
        "/Library/Fonts/SourceHanSansCN-Bold.otf",
    ),
    "Arial": (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
    ),
    "Times New Roman": (
        "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
        "/Library/Fonts/Times New Roman.ttf",
    ),
}
FALLBACK_CJK_FONT = "STHeiti"
MEASURE_FONT_CJK = "/System/Library/Fonts/Hiragino Sans GB.ttc"
MEASURE_FONT_LATIN = "/System/Library/Fonts/Supplemental/Arial.ttf"


@lru_cache(maxsize=64)
def font_available(family: str) -> bool:
    hints = FONT_FILE_HINTS.get(family, ())
    if any(Path(path).expanduser().exists() for path in hints):
        return True
    try:
        result = subprocess.run(
            ["fc-match", "-f", "%{family}\\n", family],
            check=False,
            capture_output=True,
            text=True,
            timeout=2,
        )
    except (OSError, subprocess.SubprocessError):
        return False
    families = {part.strip() for part in result.stdout.replace(",", "\n").splitlines() if part.strip()}
    return family in families


def resolve_font_family(candidates: tuple[str, ...]) -> str:
    for family in candidates:
        if font_available(family):
            return family
    return candidates[0]


DEFAULT_CJK_FONT = resolve_font_family(CJK_FONT_CANDIDATES)
DEFAULT_LATIN_FONT = resolve_font_family(LATIN_FONT_CANDIDATES)


def parse_pages(pages: str | None) -> list[int] | None:
    if not pages:
        return None
    selected: list[int] = []
    for part in pages.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start, end = int(start_s), int(end_s)
            selected.extend(range(start, end + 1))
        else:
            selected.append(int(part))
    return sorted(set(p for p in selected if p > 0))


def selected_indices(total: int, limit: int | None, pages: list[int] | None) -> list[int]:
    if pages:
        return [p - 1 for p in pages if 1 <= p <= total]
    count = total if limit is None else min(total, limit)
    return list(range(count))


def render_pdf(pdf: Path, out_dir: Path, dpi: int, limit: int | None, pages: list[int] | None) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf)
    images: list[Path] = []
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    for idx in selected_indices(len(doc), limit, pages):
        page = doc[idx]
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        path = out_dir / f"slide_{idx + 1:03d}.png"
        pix.save(path)
        images.append(path)
    return images


def natural_image_key(path: str) -> int:
    match = re.search(r"image(\d+)\.", path)
    return int(match.group(1)) if match else 10_000


def extract_pptx_images(pptx: Path, out_dir: Path, limit: int | None, pages: list[int] | None) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    images: list[Path] = []
    with zipfile.ZipFile(pptx) as zf:
        media = [
            name
            for name in zf.namelist()
            if name.startswith("ppt/media/image") and name.lower().endswith((".png", ".jpg", ".jpeg"))
        ]
        media.sort(key=natural_image_key)
        indices = selected_indices(len(media), limit, pages)
        for idx in indices:
            name = media[idx]
            suffix = Path(name).suffix.lower() or ".png"
            out = out_dir / f"slide_{idx + 1:03d}{suffix}"
            out.write_bytes(zf.read(name))
            images.append(out)
    return images


def text_color_from_region(image: Image.Image, x: int, y: int, w: int, h: int) -> str:
    # Estimate text color from the darkest non-white pixels in the OCR region.
    crop = image.crop((x, y, x + w, y + h)).convert("RGB")
    get_pixels = getattr(crop, "get_flattened_data", crop.getdata)
    pixels = list(get_pixels())
    dark = [p for p in pixels if sum(p) < 620]
    if not dark:
        return "#111111"
    avg = tuple(int(sum(c[i] for c in dark) / len(dark)) for i in range(3))
    return f"#{avg[0]:02X}{avg[1]:02X}{avg[2]:02X}"


def clean_ocr_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    # Tesseract often inserts spaces between every CJK character. Remove only
    # CJK-to-CJK spaces while keeping spaces around Latin words like GitHub.
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[，。！？：；、])", "", text)
    text = re.sub(r"(?<=[，。！？：；、])\s+(?=[\u4e00-\u9fff])", "", text)
    return text


def repair_ocr_text(text: str) -> tuple[str, list[dict[str, str]]]:
    """Apply small, explicit OCR text repairs before fusion/rebuild locking."""
    repairs: list[dict[str, str]] = []

    def replace(pattern: str, repl: str, reason: str) -> None:
        nonlocal text
        updated = re.sub(pattern, repl, text)
        if updated != text:
            repairs.append({"from": text, "to": updated, "reason": reason})
            text = updated

    replace(r"\bone[’']\s+s\b", "one’s", "english_contraction_spacing")
    replace(r"\badriver\b", "a driver", "collapsed_article_profession")
    replace(r"特殊疑问名", "特殊疑问句", "common_grammar_term_ocr_confusion")
    return text, repairs


def contains_cjk(text: str) -> bool:
    return bool(re.search(r"[\u3400-\u9fff]", text))


def suggested_font_family(text: str) -> str:
    return DEFAULT_CJK_FONT if contains_cjk(text) else DEFAULT_LATIN_FONT


def normalize_font_family(
    text: str,
    font_family: str | None,
    font_category: str | None = None,
    font_candidates: list[str] | tuple[str, ...] | None = None,
) -> str:
    category = str(font_category or "").lower()
    if contains_cjk(text):
        # Keep CJK reconstruction stable instead of inheriting mixed model/OCR
        # font guesses inside the same sentence.
        for candidate in font_candidates or []:
            if candidate in CJK_FONT_CANDIDATES and font_available(candidate):
                return candidate
        if font_family in CJK_FONT_CANDIDATES and font_available(font_family):
            return font_family
        return DEFAULT_CJK_FONT
    for candidate in font_candidates or []:
        if candidate in LATIN_FONT_CANDIDATES and font_available(candidate):
            return candidate
    if category == "serif" and font_available("Times New Roman"):
        return "Times New Roman"
    if font_family in LATIN_FONT_CANDIDATES and font_available(font_family):
        return font_family
    return DEFAULT_LATIN_FONT


def normalize_layout_fonts(layout: dict[str, Any]) -> dict[str, Any]:
    """Materialize final font families in layout JSON before saving it."""
    for page in layout.get("slides", []):
        if page.get("texts"):
            for item in page["texts"]:
                text = str(item.get("text") or item.get("textValue") or item.get("content") or "")
                family = normalize_font_family(
                    text,
                    item.get("font_family") or item.get("fontFamily"),
                    item.get("fontCategory") or item.get("font_category"),
                    item.get("fontCandidates") or item.get("font_candidates"),
                )
                item["font_family"] = family
                item["fontFamily"] = family
        for element in page.get("elements", []):
            if (element.get("type") or element.get("elementType") or "").lower() != "text":
                continue
            text = str(element.get("text") or element.get("textValue") or element.get("content") or "")
            content_data = element.get("contentData") or {}
            text = text or str(content_data.get("textValue") or content_data.get("text") or "")
            family = normalize_font_family(
                text,
                element.get("font_family") or element.get("fontFamily"),
                element.get("fontCategory") or element.get("font_category"),
                element.get("fontCandidates") or element.get("font_candidates"),
            )
            element["fontFamily"] = family
            element["font_family"] = family
    return layout


def sanitize_pptx_xml_fonts(xml: str) -> str:
    """Remove template font noise so the deck advertises only approved fonts."""
    replacements = {
        "+mn-lt": DEFAULT_LATIN_FONT,
        "+mj-lt": DEFAULT_LATIN_FONT,
        "+mn-cs": DEFAULT_LATIN_FONT,
        "+mj-cs": DEFAULT_LATIN_FONT,
        "+mn-ea": DEFAULT_CJK_FONT,
        "+mj-ea": DEFAULT_CJK_FONT,
    }
    for old, new in replacements.items():
        xml = xml.replace(f'typeface="{old}"', f'typeface="{new}"')

    xml = re.sub(r'(<a:latin\b[^>]*\btypeface=")[^"]*(")', rf"\g<1>{DEFAULT_LATIN_FONT}\2", xml)
    xml = re.sub(r'(<a:ea\b[^>]*\btypeface=")[^"]*(")', rf"\g<1>{DEFAULT_CJK_FONT}\2", xml)
    xml = re.sub(r'(<a:cs\b[^>]*\btypeface=")[^"]*(")', rf"\g<1>{DEFAULT_LATIN_FONT}\2", xml)
    xml = re.sub(r'(<a:buFont\b[^>]*\btypeface=")[^"]*(")', rf"\g<1>{DEFAULT_LATIN_FONT}\2", xml)

    cjk_scripts = {"Hans", "Hant", "Jpan", "Kore"}
    return re.sub(r'(<a:font\b[^>]*\bscript="([^"]+)"[^>]*\btypeface=")[^"]*(")', lambda m: f'{m.group(1)}{DEFAULT_CJK_FONT if m.group(2) in cjk_scripts else DEFAULT_LATIN_FONT}{m.group(3)}', xml)


def sanitize_pptx_fonts(pptx_path: Path) -> None:
    """Rewrite non-slide theme/master font declarations after python-pptx save."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp_path = Path(tmp.name)
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if (
                    info.filename.endswith(".xml")
                    and not info.filename.startswith("ppt/slides/")
                    and info.filename.startswith("ppt/")
                ):
                    text = data.decode("utf-8", errors="ignore")
                    data = sanitize_pptx_xml_fonts(text).encode("utf-8")
                zout.writestr(info, data)
        shutil.move(str(tmp_path), pptx_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink()


def measure_font(text: str, size_px: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [MEASURE_FONT_CJK, MEASURE_FONT_LATIN] if contains_cjk(text) else [MEASURE_FONT_LATIN, MEASURE_FONT_CJK]
    for path in candidates:
        try:
            return ImageFont.truetype(path, max(1, int(size_px)))
        except OSError:
            continue
    return ImageFont.load_default()


def text_bbox_size(text: str, size_px: float) -> tuple[float, float]:
    font = measure_font(text, int(size_px))
    lines = str(text).splitlines() or [str(text)]
    widths: list[float] = []
    heights: list[float] = []
    for line in lines:
        bbox = font.getbbox(line or " ")
        widths.append(float(bbox[2] - bbox[0]))
        heights.append(float(bbox[3] - bbox[1]))
    return max(widths or [0.0]), sum(heights or [0.0]) + max(0, len(lines) - 1) * size_px * 0.25


def fit_font_size_px(text: str, box_w: float, box_h: float, requested_px: float) -> float:
    """Find the largest practical source-coordinate font size that fits the box."""
    requested = max(8.0, float(requested_px))
    upper = min(max(requested, float(box_h) * 0.92), requested * 1.45, 140.0)
    size = max(8.0, upper)
    max_w = max(float(box_w) * 0.98, 8.0)
    max_h = max(float(box_h) * 1.18, 8.0)
    for _ in range(32):
        measured_w, measured_h = text_bbox_size(text, size)
        if measured_w <= max_w and measured_h <= max_h:
            return size
        ratio_w = max_w / measured_w if measured_w else 1.0
        ratio_h = max_h / measured_h if measured_h else 1.0
        size *= max(0.55, min(ratio_w, ratio_h, 0.96))
        if size <= 8:
            return 8.0
    return max(8.0, size)


def useful_text(text: str, conf: float, width: int, height: int) -> bool:
    if not text:
        return False
    alnum_or_cjk = re.findall(r"[A-Za-z0-9\u4e00-\u9fff]", text)
    if len(alnum_or_cjk) < 4:
        return False
    if conf < 45 and len(alnum_or_cjk) < 5:
        return False
    if width < 25 or height < 12:
        return False
    return True


def filter_line_word_outliers(line_words: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Remove obvious OCR word outliers before forming an editable text line."""
    if len(line_words) < 4:
        return line_words
    heights = sorted(float(w.get("height") or 0) for w in line_words if float(w.get("height") or 0) > 0)
    if not heights:
        return line_words
    median_h = heights[len(heights) // 2]
    if median_h <= 0:
        return line_words
    kept: list[dict[str, Any]] = []
    for word in line_words:
        h = float(word.get("height") or 0)
        text = str(word.get("text") or "").strip()
        # Large isolated glyphs from illustrations or background labels can be
        # merged into a real line by Tesseract. They should not determine text
        # content, bbox, or font size for the editable OCR line.
        if h > median_h * 1.55 and len(text) <= 2:
            continue
        kept.append(word)
    return kept or line_words


def ocr_slide(image_path: Path, lang: str, min_conf: int, psm: int) -> dict[str, Any]:
    image = Image.open(image_path).convert("RGB")
    data = pytesseract.image_to_data(
        image,
        lang=lang,
        config=f"--psm {psm}",
        output_type=pytesseract.Output.DICT,
    )
    words: list[dict[str, Any]] = []
    mask_words: list[dict[str, Any]] = []
    for idx, text in enumerate(data["text"]):
        text = (text or "").strip()
        if not text:
            continue
        try:
            conf = float(data["conf"][idx])
        except ValueError:
            conf = -1
        x, y, w, h = (int(data[k][idx]) for k in ("left", "top", "width", "height"))
        cleaned_word = clean_ocr_text(text)
        if not cleaned_word or w < 8 or h < 8:
            continue
        item = {
            "text": cleaned_word,
            "confidence": conf,
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "block": data["block_num"][idx],
            "paragraph": data["par_num"][idx],
            "line": data["line_num"][idx],
            "color": text_color_from_region(image, x, y, w, h),
        }
        # Masking should be permissive so low-confidence but visibly real text
        # gets removed from the background. Editable overlay text remains
        # stricter via the grouped-line filter below.
        if conf >= 0:
            mask_words.append(item)
        if conf >= min_conf:
            words.append(item)

    lines: list[dict[str, Any]] = []
    grouped: dict[tuple[int, int, int], list[dict[str, Any]]] = {}
    for word in words:
        key = (word["block"], word["paragraph"], word["line"])
        grouped.setdefault(key, []).append(word)
    for key, line_words in grouped.items():
        line_words.sort(key=lambda w: w["x"])
        line_words = filter_line_word_outliers(line_words)
        raw_line_text = clean_ocr_text(" ".join(w["text"] for w in line_words))
        text, text_repairs = repair_ocr_text(raw_line_text)
        x1 = min(w["x"] for w in line_words)
        y1 = min(w["y"] for w in line_words)
        x2 = max(w["x"] + w["width"] for w in line_words)
        y2 = max(w["y"] + w["height"] for w in line_words)
        avg_conf = sum(w["confidence"] for w in line_words) / len(line_words)
        if not useful_text(text, avg_conf, x2 - x1, y2 - y1):
            continue
        font_px = max(w["height"] for w in line_words)
        lines.append(
            {
                "text": text,
                "raw_text": raw_line_text,
                "text_repairs": text_repairs,
                "confidence": round(avg_conf, 1),
                "x": x1,
                "y": y1,
                "width": x2 - x1,
                "height": y2 - y1,
                "font_size_px": font_px,
                "font_family": suggested_font_family(text),
                "color": line_words[0]["color"],
                "fill_color": "#{:02X}{:02X}{:02X}".format(*sample_fill_color(image, (x1, y1, x2, y2))),
                "key": key,
                "word_boxes": [
                    {
                        "x": w["x"],
                        "y": w["y"],
                        "width": w["width"],
                        "height": w["height"],
                        "text": w["text"],
                    }
                    for w in line_words
                ],
            }
        )

    return {
        "image": str(image_path),
        "width": image.width,
        "height": image.height,
        "words": mask_words,
        "texts": lines,
    }


def hex_to_rgb(color: str) -> RGBColor | None:
    color = color.lstrip("#")
    if len(color) != 6:
        return None
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def rgb_to_hex(rgb: list[int] | tuple[int, int, int] | None, default: str = "#111111") -> str:
    if not rgb or len(rgb) < 3:
        return default
    return "#{:02X}{:02X}{:02X}".format(
        max(0, min(255, int(rgb[0]))),
        max(0, min(255, int(rgb[1]))),
        max(0, min(255, int(rgb[2]))),
    )


def normalize_text_item(item: dict[str, Any]) -> dict[str, Any]:
    """Normalize OCR, repaired, and Codia-like text fields into one shape."""
    text = item.get("text") or item.get("textValue") or item.get("content") or ""
    content_data = item.get("contentData") or {}
    if not text:
        text = content_data.get("textValue") or content_data.get("text") or ""

    style = item.get("styleConfig") or item.get("style") or {}
    text_config = item.get("textConfig") or style.get("textConfig") or {}
    layout = item.get("layoutConfig") or item.get("layout") or {}
    absolute = layout.get("absoluteAttrs") or {}
    coord = item.get("coord") or absolute.get("coord") or [item.get("x", 0), item.get("y", 0)]
    width_spec = style.get("widthSpec") or {}
    height_spec = style.get("heightSpec") or {}

    color = item.get("color")
    if isinstance(color, list):
        color = rgb_to_hex(color)
    if not color:
        text_color = style.get("textColor") or item.get("textColor") or {}
        color = rgb_to_hex(text_color.get("rgbValues"), "#111111")

    font_size = (
        item.get("font_size_px")
        or item.get("fontSize")
        or text_config.get("fontSize")
        or max(12, item.get("height", height_spec.get("value", 24)) * 0.7)
    )
    font_family = item.get("font_family") or item.get("fontFamily") or text_config.get("fontFamily") or suggested_font_family(text)
    align = item.get("align") or item.get("textAlign") or text_config.get("textAlign") or ["LEFT", "CENTER"]

    font_size_locked = bool(item.get("fontSizeLocked") or item.get("font_size_locked"))
    position_locked = bool(item.get("positionLocked") or item.get("position_locked"))
    text_source = item.get("textSource") or item.get("text_source") or ""
    preserve_text = bool(item.get("matchTrusted") or text_source == "ocr_exact" or (font_size_locked and position_locked))
    raw_text = str(text)
    if preserve_text:
        text = raw_text
    elif "\n" in raw_text:
        text = "\n".join(clean_ocr_text(line) for line in raw_text.splitlines())
    else:
        text = clean_ocr_text(raw_text)
    if font_size_locked:
        fitted_font_size = float(font_size)
    else:
        fitted_font_size = fit_font_size_px(
            text,
            float(item.get("width", width_spec.get("value", 100))),
            float(item.get("height", height_spec.get("value", max(12, float(font_size))))),
            float(font_size),
        )

    return {
        "text": text,
        "x": float(item.get("x", coord[0] if len(coord) > 0 else 0)),
        "y": float(item.get("y", coord[1] if len(coord) > 1 else 0)),
        "width": float(item.get("width", width_spec.get("value", 100))),
        "height": float(item.get("height", height_spec.get("value", max(12, float(font_size))))),
        "font_size_px": fitted_font_size,
        "font_size_locked": font_size_locked,
        "position_locked": position_locked,
        "fontSizeSource": item.get("fontSizeSource") or item.get("font_size_source") or "unknown",
        "positionSource": item.get("positionSource") or item.get("position_source") or "unknown",
        "textSource": text_source or item.get("textSource") or item.get("text_source") or "unknown",
        "font_family": normalize_font_family(
            text,
            font_family,
            item.get("fontCategory") or item.get("font_category"),
            item.get("fontCandidates") or item.get("font_candidates"),
        ),
        "color": color,
        "fill_color": item.get("fill_color") or item.get("fillColor") or "#FFFFFF",
        "align": align,
        "bold": item.get("bold") or str(text_config.get("fontStyle", "")).lower() in {"bold", "extra_bold", "semi_bold"},
        "fontCategory": item.get("fontCategory") or item.get("font_category") or "unknown",
        "fontCandidates": item.get("fontCandidates") or item.get("font_candidates") or [],
        "fontWeight": item.get("fontWeight") or item.get("font_weight") or (700 if item.get("bold") else 400),
        "fontWeightSource": item.get("fontWeightSource") or item.get("font_weight_source") or "unknown",
        "fontWeightLocked": bool(item.get("fontWeightLocked") or item.get("font_weight_locked")),
        "fontWeightConfidence": item.get("fontWeightConfidence") or item.get("font_weight_confidence") or 0,
        "line_spacing": item.get("lineSpacing") or item.get("line_spacing"),
        "styleEvidence": item.get("styleEvidence") or item.get("style_evidence") or [],
        "styleSource": item.get("styleSource") or item.get("style_source") or "unknown",
        "styleConfidence": item.get("styleConfidence") or item.get("style_confidence") or 0,
    }


def page_text_items(page: dict[str, Any]) -> list[dict[str, Any]]:
    if page.get("texts"):
        return [normalize_text_item(item) for item in page["texts"]]
    items: list[dict[str, Any]] = []
    for element in page.get("elements", []):
        if (element.get("type") or element.get("elementType") or "").lower() == "text":
            items.append(normalize_text_item(element))
    return items


def page_image_items(page: dict[str, Any]) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for element in page.get("elements", []):
        if (element.get("type") or element.get("elementType") or "").lower() != "image":
            continue
        src = element.get("src") or element.get("image") or (element.get("contentData") or {}).get("imageSource")
        if not src:
            continue
        if str(src) == str(page.get("image")) or str(src) == str(page.get("clean_background")):
            continue
        if not (str(src).startswith("http://") or str(src).startswith("https://")) and not Path(str(src)).expanduser().exists():
            continue
        items.append(
            {
                "src": src,
                "x": float(element.get("x", 0)),
                "y": float(element.get("y", 0)),
                "width": float(element.get("width", 100)),
                "height": float(element.get("height", 100)),
            }
        )
    return items


def page_shape_items(page: dict[str, Any]) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for element in page.get("elements", []):
        if (element.get("type") or element.get("elementType") or "").lower() != "shape":
            continue
        items.append(
            {
                "x": float(element.get("x", 0)),
                "y": float(element.get("y", 0)),
                "width": float(element.get("width", 0)),
                "height": float(element.get("height", 0)),
                "shape": element.get("shape"),
                "fill_color": element.get("fill_color") or element.get("fillColor") or "#FFFFFF",
            }
        )
    return items


def add_shape_element(slide, item: dict[str, Any], img_w: int, img_h: int, slide_w: float, slide_h: float) -> None:
    x = item["x"] / img_w * slide_w
    y = item["y"] / img_h * slide_h
    w = max(item["width"] / img_w * slide_w, 0.01)
    h = max(item["height"] / img_h * slide_h, 0.01)
    shape_kind = str(item.get("shape") or "").lower()
    mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE if shape_kind == "rounded_rect" else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(mso_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    rgb = hex_to_rgb(item.get("fill_color", "#FFFFFF")) or RGBColor(255, 255, 255)
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_textbox(slide, item: dict[str, Any], img_w: int, img_h: int, slide_w: float, slide_h: float, fill_background: bool = False) -> None:
    item = normalize_text_item(item)
    x = item["x"] / img_w * slide_w
    y = item["y"] / img_h * slide_h
    w = max(item["width"] / img_w * slide_w, 0.2)
    height_scale = 1.0 if item.get("position_locked") else 2.0
    h = max(item["height"] / img_h * slide_h * height_scale, 0.16)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_background:
        fill = box.fill
        fill.solid()
        rgb = hex_to_rgb(item.get("fill_color", "#FFFFFF")) or RGBColor(255, 255, 255)
        fill.fore_color.rgb = rgb
        box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    align = item.get("align", ["LEFT", "CENTER"])
    if isinstance(align, list):
        align = align[0] if align else "LEFT"
    def style_paragraph(paragraph) -> None:
        paragraph.alignment = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFIED": PP_ALIGN.JUSTIFY,
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(str(align), PP_ALIGN.LEFT)
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = float(item.get("line_spacing") or 0.92)

    style_paragraph(p)
    font_name = item.get("font_family") or (DEFAULT_CJK_FONT if contains_cjk(item["text"]) else DEFAULT_LATIN_FONT)
    font_size = Pt(max(1, item["font_size_px"] * slide_h / img_h * 72 * PPT_RENDER_FONT_SCALE))
    color = item.get("color", "#111111").lstrip("#")

    def style_run(run) -> None:
        font = run.font
        font.name = font_name
        rpr = run._r.get_or_add_rPr()
        latin = rpr.get_or_add_latin()
        latin.set("typeface", font_name)
        for tag in ("a:ea", "a:cs"):
            node = rpr.find(qn(tag))
            if node is None:
                node = OxmlElement(tag)
                rpr.append(node)
            node.set("typeface", font_name)
        font.size = font_size
        font.bold = bool(item.get("bold"))
        if len(color) == 6:
            font.color.rgb = RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))

    lines = str(item["text"]).splitlines() or [""]
    run = p.add_run()
    run.text = lines[0]
    style_run(run)
    for line in lines[1:]:
        p = tf.add_paragraph()
        style_paragraph(p)
        run = p.add_run()
        run.text = line
        style_run(run)


def sample_fill_color(image: Image.Image, box: tuple[int, int, int, int]) -> tuple[int, int, int]:
    x1, y1, x2, y2 = box
    w, h = image.size
    pad = 8
    samples: list[tuple[int, int, int]] = []
    regions = [
        (max(0, x1 - pad), max(0, y1 - pad), min(w, x2 + pad), max(0, y1)),
        (max(0, x1 - pad), min(h, y2), min(w, x2 + pad), min(h, y2 + pad)),
        (max(0, x1 - pad), max(0, y1), max(0, x1), min(h, y2)),
        (min(w, x2), max(0, y1), min(w, x2 + pad), min(h, y2)),
    ]
    for region in regions:
        if region[2] <= region[0] or region[3] <= region[1]:
            continue
        crop = image.crop(region).convert("RGB")
        get_pixels = getattr(crop, "get_flattened_data", crop.getdata)
        samples.extend(list(get_pixels()))
    if not samples:
        return (255, 255, 255)
    # Use bright/background-biased pixels so dark text does not pollute the fill.
    bright = [p for p in samples if sum(p) > 560]
    pool = bright or samples
    return tuple(int(sum(p[i] for p in pool) / len(pool)) for i in range(3))


def create_clean_background(page: dict[str, Any], out_dir: Path, expand_px: int) -> str:
    source = Path(page["image"])
    image = Image.open(source).convert("RGB")
    mask = Image.new("L", image.size, 0)
    mask_draw = ImageDraw.Draw(mask)
    draw = ImageDraw.Draw(image)
    mask_items: list[dict[str, Any]] = []
    text_lines = page.get("texts")
    if text_lines:
        for text_line in text_lines:
            mask_items.extend(text_line.get("word_boxes") or [text_line])
    else:
        mask_items.extend(page_text_items(page))
    mask_items.extend(page.get("mask_texts") or [])
    for item in mask_items:
        raw_x1 = float(item.get("x", 0))
        raw_y1 = float(item.get("y", 0))
        raw_x2 = raw_x1 + float(item.get("width", 0))
        raw_y2 = raw_y1 + float(item.get("height", 0))
        raw_x1, raw_x2 = sorted((raw_x1, raw_x2))
        raw_y1, raw_y2 = sorted((raw_y1, raw_y2))
        x1 = max(0, int(raw_x1) - expand_px)
        y1 = max(0, int(raw_y1) - expand_px)
        x2 = min(image.width, int(raw_x2) + expand_px)
        y2 = min(image.height, int(raw_y2) + expand_px)
        if x2 <= x1 or y2 <= y1:
            continue
        fill = sample_fill_color(image, (x1, y1, x2, y2))
        draw.rounded_rectangle((x1, y1, x2, y2), radius=max(2, expand_px // 2), fill=fill)
        mask_draw.rectangle((x1, y1, x2, y2), fill=255)
    clean_dir = out_dir / "cleaned"
    mask_dir = out_dir / "masks"
    clean_dir.mkdir(parents=True, exist_ok=True)
    mask_dir.mkdir(parents=True, exist_ok=True)
    clean_path = clean_dir / source.name
    mask_path = mask_dir / source.with_suffix(".mask.png").name
    image.save(clean_path)
    mask.save(mask_path)
    page["clean_background"] = str(clean_path)
    page["text_mask"] = str(mask_path)
    return str(clean_path)


def add_image_element(slide, item: dict[str, Any], img_w: int, img_h: int, slide_w: float, slide_h: float) -> None:
    src = item["src"]
    if src.startswith("http://") or src.startswith("https://"):
        # The local rebuild flow should use downloaded assets. Keep remote
        # images in layout JSON, but do not fetch them implicitly here.
        return
    x = item["x"] / img_w * slide_w
    y = item["y"] / img_h * slide_h
    w = item["width"] / img_w * slide_w
    h = item["height"] / img_h * slide_h
    slide.shapes.add_picture(src, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def build_pptx(layout: dict[str, Any], output: Path, background_key: str = "image", overlay_opacity_note: bool = True, fill_text_background: bool = False) -> None:
    prs = Presentation()
    prs.slide_width = int(DEFAULT_SLIDE_W * EMU_PER_INCH)
    prs.slide_height = int(DEFAULT_SLIDE_H * EMU_PER_INCH)
    blank = prs.slide_layouts[6]
    for page in layout["slides"]:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            page.get(background_key) or page["image"],
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        for item in page_shape_items(page):
            add_shape_element(slide, item, page["width"], page["height"], DEFAULT_SLIDE_W, DEFAULT_SLIDE_H)
        for item in page_image_items(page):
            add_image_element(slide, item, page["width"], page["height"], DEFAULT_SLIDE_W, DEFAULT_SLIDE_H)
        for item in page_text_items(page):
            add_textbox(slide, item, page["width"], page["height"], DEFAULT_SLIDE_W, DEFAULT_SLIDE_H, fill_background=fill_text_background)
        if overlay_opacity_note:
            slide.notes_slide.notes_text_frame.text = (
                f"Prototype {layout.get('mode')} slide: background plus layout-derived editable elements."
            )
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    sanitize_pptx_fonts(output)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create editable text-overlay deck from image-based slides")
    source = parser.add_mutually_exclusive_group(required=True)
    source.add_argument("--pdf")
    source.add_argument("--pptx")
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--limit", type=int)
    parser.add_argument("--pages", help="Comma-separated 1-based pages or ranges, e.g. 1,2,5-7")
    parser.add_argument("--dpi", type=int, default=180)
    parser.add_argument("--lang", default="chi_sim+eng")
    parser.add_argument("--min-conf", type=int, default=35)
    parser.add_argument("--psm", type=int, default=11)
    parser.add_argument("--mode", choices=["text-overlay", "clean-text", "cover-text"], default="text-overlay")
    parser.add_argument("--background-key", default=None, help="Use this page key as slide background, e.g. clean_background")
    parser.add_argument("--mask-expand", type=int, default=6)
    parser.add_argument("--write-raw", action="store_true", help="Also write layout.raw.json before mode-specific mutations")
    parser.add_argument("--repair-layout", help="Reserved path for a future model-repaired layout JSON")
    args = parser.parse_args()

    out_dir = Path(args.output_dir).expanduser().resolve()
    rendered_dir = out_dir / "rendered"
    pptx_path = out_dir / "editable_text_overlay.pptx"
    layout_path = out_dir / "layout.json"

    pages = parse_pages(args.pages)
    if args.pdf:
        source_path = Path(args.pdf).expanduser().resolve()
        image_paths = render_pdf(source_path, rendered_dir, args.dpi, args.limit, pages)
    else:
        source_path = Path(args.pptx).expanduser().resolve()
        image_paths = extract_pptx_images(source_path, rendered_dir, args.limit, pages)
    slides = [ocr_slide(path, args.lang, args.min_conf, args.psm) for path in image_paths]
    layout = {
        "source": str(source_path),
        "mode": args.mode,
        "dpi": args.dpi,
        "pages": pages,
        "slides": slides,
    }
    raw_layout_path = out_dir / "layout.raw.json"
    if args.write_raw:
        raw_layout_path.write_text(json.dumps(layout, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.repair_layout:
        repaired_path = Path(args.repair_layout).expanduser().resolve()
        if repaired_path.exists():
            layout = json.loads(repaired_path.read_text(encoding="utf-8"))
            layout["mode"] = args.mode
        else:
            raise FileNotFoundError(f"repair layout not found: {repaired_path}")
    background_key = args.background_key or "image"
    if args.mode == "clean-text":
        for page in layout["slides"]:
            create_clean_background(page, out_dir, args.mask_expand)
        background_key = args.background_key or "clean_background"
    normalize_layout_fonts(layout)
    layout_path.write_text(json.dumps(layout, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.mode == "clean-text":
        pptx_path = out_dir / "editable_clean_text.pptx"
    elif args.mode == "cover-text":
        pptx_path = out_dir / "editable_cover_text.pptx"
    build_pptx(layout, pptx_path, background_key=background_key, fill_text_background=args.mode == "cover-text")
    print(json.dumps({
        "ok": True,
        "slides": len(layout["slides"]),
        "pages": pages,
        "text_nodes": sum(len(page_text_items(page)) for page in layout["slides"]),
        "image_nodes": sum(len(page_image_items(page)) for page in layout["slides"]),
        "shape_nodes": sum(len(page_shape_items(page)) for page in layout["slides"]),
        "raw_layout": str(raw_layout_path) if args.write_raw else None,
        "layout": str(layout_path),
        "pptx": str(pptx_path),
    }, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
